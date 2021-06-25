# -*- coding: utf-8 -*-
"""
@Author:
  Samuel Frederick, NSSC Contractor (ORAU)
  U.S. EPA, Office of Research and Development
  Center for Environmental Measurement and Modeling
  Air Methods and Characterization Division, Source and Fine Scale Branch
  109 T.W Alexander Drive, Research Triangle Park, NC 27711
  Office: 919-541-4086 | Email: frederick.samuel@epa.gov

Created:
  Tue Dec 15 08:53:19 2020
Last Updated:
  Thu Jun 24 11:24:00 2021
"""

import pptx as ppt
import datetime as dt
import numpy as np
import math
import os
import sys
lib_path = os.path.abspath(__file__ + '../../..')
if lib_path not in sys.path:
    sys.path.append(lib_path)
from Sensor_Evaluation.sensor_eval_class import SensorEvaluation


class PerformanceReport(SensorEvaluation):
    """
    Generate air sensor performance reporting document for evaluations
    conducted under fixed site, ambient, outdoor settings.

    Configured currently to create reports for PM2.5 and O3.
    """

    def __init__(self, sensor_name, eval_param, load_raw_data=False,
                 reference_data=None, ref_name=None,
                 serials=None, tzone_shift=0, bbox=None, aqs_id=None,
                 write_to_file=False, fmt_sensor_name=None, testing_org=None,
                 testing_loc=None):

        # Inherit the SensorEvaluation class instance attributes
        super().__init__(sensor_name, eval_param, load_raw_data,
                         reference_data, ref_name, serials, tzone_shift,
                         bbox, aqs_id, write_to_file)

        self.fmt_sensor_name = fmt_sensor_name
        # Placeholder method for formatted sensor name, replace '_' with spaces
        if self.fmt_sensor_name is None:
            self.fmt_sensor_name = self.sensor_name.replace('_', ' ')
        self.today = dt.datetime.now().strftime('%y%m%d')

        self.template_name = ('Reporting_Template_Base_' + self.eval_param
                              + '.pptx')
        # Path to reporting template
        self.template_path = '\\'.join((self.lib_path, 'Reports', 'templates',
                                        self.eval_param, self.template_name))

        # Details about testing and deployment site
        self.testing_org = testing_org
        self.testing_loc = testing_loc

        # Populate deployment dictionary with performance metric results
        self.calculate_metrics()

        # Sampling timeframe
        self.tframe = {grp:
                       self.deploy_dict['Deployment Groups'][grp]['eval_start']
                       + ' - ' +
                       self.deploy_dict['Deployment Groups'][grp]['eval_end']
                       for grp in list(
                               self.deploy_dict['Deployment Groups'].keys())
                       }

        # Keys are sensor serial IDs, values are deployment group number
        # Useful if multiple evaluation groups deployed
        self.serial_grp_dict = {}
        for grp in self.deploy_dict['Deployment Groups']:
            grp_dict = self.deploy_dict['Deployment Groups'][grp]
            grp_sensors = grp_dict['sensors']
            for sensor in grp_sensors:
                serial = grp_sensors[sensor]['serial_id']
                self.serial_grp_dict[serial] = grp

        # Initialize report object
        self.rpt = ppt.Presentation(self.template_path)
        self.shapes = self.rpt.slides[0].shapes

        # Shape at backgroud around which to orient other figures
        self.cursor_sp = self.shapes[0]._element

        # Initialize figure positions in report
        self.FigPositions()

    def FigPositions(self):
        """
        Figure positions for reports. Values are in inches, specifying the
        left and top center location of each figure.
        """
        self.fig_locs = {
                        'SingleScatter': {'left': '',
                                          'top': ''},
                        'TripleScatter': {'left': '',
                                          'top': ''},
                        'Timeseries': {'left': '',
                                       'top': ''},
                        'MetricPlot': {'left': '',
                                       'top': ''},
                        'MetDist': {'left': '',
                                    'top': ''},
                        'MetInfl': {'left': '',
                                    'top': ''}
                       }

        if self.eval_param == 'PM25':
            self.fig_locs['SingleScatter']['left'] = 11.11
            self.fig_locs['SingleScatter']['top'] = 8.16
            self.fig_locs['TripleScatter']['left'] = 2.35
            self.fig_locs['TripleScatter']['top'] = 3.82
            self.fig_locs['Timeseries']['left'] = 0.63
            self.fig_locs['Timeseries']['top'] = 8.15
            self.fig_locs['MetricPlot']['left'] = 0.65
            self.fig_locs['MetricPlot']['top'] = 13.19
            self.fig_locs['MetDist']['left'] = 1.91
            self.fig_locs['MetDist']['top'] = 17.58
            self.fig_locs['MetInfl']['left'] = 8.24
            self.fig_locs['MetInfl']['top'] = 17.54

        if self.eval_param == 'O3':
            self.fig_locs['SingleScatter']['left'] = 11.57
            self.fig_locs['SingleScatter']['top'] = 8.14
            self.fig_locs['TripleScatter']['left'] = 2.35
            self.fig_locs['TripleScatter']['top'] = 3.82
            self.fig_locs['Timeseries']['left'] = 0.63
            self.fig_locs['Timeseries']['top'] = 8.19
            self.fig_locs['MetricPlot']['left'] = 0.7
            self.fig_locs['MetricPlot']['top'] = 12.78
            self.fig_locs['MetDist']['left'] = 1.6
            self.fig_locs['MetDist']['top'] = 17.3
            self.fig_locs['MetInfl']['left'] = 8.28
            self.fig_locs['MetInfl']['top'] = 17.31

    def AddSingleScatterPlot(self, **kwargs):
        """
        Add sensor vs. reference scatter plots (1-hr, 24-hr [PM2.5 only])
        to report
        """
        fig_name = self.sensor_name + '_vs_' + self.ref_name + '_report_fmt'

        # Search for figure created today
        try:
            fig_name += '_' + self.today + '.png'
            fig_path = self.figure_path + '\\'.join((self.eval_param,
                                                     fig_name))
            figure = open(fig_path, 'r')
            figure.close()

        # If figure not found, load sensor data and create figure
        except FileNotFoundError:

            self.plot_sensor_scatter(
                        plot_subset=kwargs.get('plot_subset', ['1']),
                        plot_limits=kwargs.get('plot_limits', (-1, 30)),
                        tick_spacing=kwargs.get('tick_spacing', 5),
                        text_pos=kwargs.get('text_pos', 'upper_left'),
                        report_fmt=True)

        if len(kwargs) != 0:
            self.plot_sensor_scatter(
                    plot_subset=kwargs.get('plot_subset', ['1']),
                    plot_limits=kwargs.get('plot_limits', (-1, 30)),
                    tick_spacing=kwargs.get('tick_spacing', 5),
                    text_pos=kwargs.get('text_pos', 'upper_left'),
                    report_fmt=True)

        scatter_loc = self.fig_locs['SingleScatter']
        self.scatterplt = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(scatter_loc['left']),
                                    top=ppt.util.Inches(scatter_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.scatterplt._element)

    def AddMultiScatter(self, **kwargs):
        """
        Add sensor vs. reference scatter plots 24-hr for all sensors in testing
        group.
        """
        # Use slide layout for generating additional slides
        slide_layout_idx = 0
        slide_layout = self.rpt.slide_layouts[slide_layout_idx]
        # Create new page
        slide = self.rpt.slides.add_slide(slide_layout)

        # Add Sensor-Sensor section header text label
        section_header = slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(2.81),  # top
                                ppt.util.Inches(3.16),  # width
                                ppt.util.Inches(0.47))  # height
        section_header_obj = section_header.text_frame.paragraphs[0]
        section_header_obj.text = 'Sensor-Reference Scatter Plots'
        self.FormatText(section_header_obj, alignment='left',
                        font_name='Calibri Light', font_size=22)

        # A horizontal line separating the header text from the section body
        hline = slide.shapes.add_connector(
                            ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                            ppt.util.Inches(0.85),  # start_x
                            ppt.util.Inches(3.30),  # start_y
                            ppt.util.Inches(16.24),  # end_x
                            ppt.util.Inches(3.30))  # end_y

        hline.line.fill.solid()
        hline.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171, 171, 171)

        # Loop over averaging intervals specified for the parameter
        for i, avg_interval in enumerate(self.eval_param_averaging):
            fig_name = (self.sensor_name + '_vs_' + self.ref_name +
                        '_' + avg_interval + '_3_sensors')

            # Search for figure created today
            try:
                fig_name += '_' + self.today + '.png'
                fig_path = self.figure_path + '\\' + self.eval_param + '\\' \
                    + fig_name

                figure = open(fig_path, 'r')
                figure.close()

            # If figure not found, load sensor data and create figure
            except FileNotFoundError as e:
                print(e)

                if len(kwargs) == 0:
                    print('Warning: No plotting arguments passed to function,'
                          ' using default configuration')
                self.plot_sensor_scatter(
                            avg_interval,
                            plot_limits=kwargs.get('plot_limits', (-1, 30)),
                            tick_spacing=kwargs.get('tick_spacing', 5),
                            text_pos=kwargs.get('text_pos', 'upper_left'))

            print(fig_path)

            if len(kwargs) != 0:
                self.plot_sensor_scatter(
                        avg_interval,
                        plot_limits=kwargs.get('plot_limits', (-1, 30)),
                        tick_spacing=kwargs.get('tick_spacing', 5),
                        text_pos=kwargs.get('text_pos', 'upper_left'))

            if self.n_sensors <= 3:
                scatter_loc = self.fig_locs['TripleScatter']
                fig_height = 5.62  # height of triple scatter figure in inches
            else:
                # Add a page, place figure on new page
                scatter_loc = self.fig_locs['MultiScatter']

                # TODO: set correct figure height for figures with mult. rows
                fig_height = 5.62  # height of triple scatter figure in inches

            left = ppt.util.Inches(scatter_loc['left'])
            top = ppt.util.Inches(scatter_loc['top']+ i*fig_height)

            slide.shapes.add_picture(fig_path, left, top)

    def AddTimeseriesPlot(self, **kwargs):
        """
        Add timeseries plots (1-hr, 24-hr [PM2.5 only]) to report
        """
        fig_name = self.sensor_name + '_timeseries_' + self.eval_param \
            + '_report_fmt'

        # Search for figure created today
        try:
            fig_name += '_' + self.today + '.png'
            fig_path = self.figure_path + '\\'.join((self.eval_param,
                                                     fig_name))

            figure = open(fig_path, 'r')
            figure.close()

        except FileNotFoundError:
            self.plot_timeseries(
                    format_xaxis_weeks=kwargs.get('format_xaxis_weeks', False),
                    yscale=kwargs.get('yscale', 'linear'),
                    date_interval=kwargs.get('date_interval', 7),
                    report_fmt=True,
                    ylims=kwargs.get('ylims', (0, 25)))

        if len(kwargs) != 0:
            self.plot_timeseries(
                format_xaxis_weeks=kwargs.get('format_xaxis_weeks', False),
                yscale=kwargs.get('yscale', 'linear'),
                date_interval=kwargs.get('date_interval', 7),
                report_fmt=True,
                ylims=kwargs.get('ylims', (0, 25)))

        timeseries_loc = self.fig_locs['Timeseries']
        self.timeseries = self.shapes.add_picture(
                                fig_path,
                                left=ppt.util.Inches(timeseries_loc['left']),
                                top=ppt.util.Inches(timeseries_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.timeseries._element)

    def AddMetricsPlot(self, **kwargs):
        """
        Add Performance target metric boxplots/dot plots to report,
        (1-hr, 24-hr [PM2.5 only])
        """
        fig_name = self.sensor_name + '_regression_boxplot_' + self.eval_param

        # Search for figure created today
        try:
            fig_name += '_' + self.today + '.png'
            fig_path = self.figure_path + '\\'.join((self.eval_param,
                                                     fig_name))

            figure = open(fig_path, 'r')
            figure.close()

        except FileNotFoundError:

            self.plot_metrics()

        metricplt_loc = self.fig_locs['MetricPlot']
        self.metricplot = self.shapes.add_picture(
                                fig_path,
                                left=ppt.util.Inches(metricplt_loc['left']),
                                top=ppt.util.Inches(metricplt_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metricplot._element)

    def AddMetDistPlot(self, **kwargs):
        """
        Add meteorological distribution (temperature, relative humidity) to
        report
        """
        fig_name = self.sensor_name + '_met_distplot_report_fmt'

        # Search for figure created today
        try:
            fig_name += '_' + self.today + '.png'
            fig_path = self.figure_path + '\\'.join(('Met', fig_name))

            figure = open(fig_path, 'r')
            figure.close()

        except FileNotFoundError:

            self.plot_met_dist()

        metdist_loc = self.fig_locs['MetDist']
        self.metdist = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(metdist_loc['left']),
                                    top=ppt.util.Inches(metdist_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metdist._element)

    def AddMetInflPlot(self, **kwargs):
        """
        Add normalized meteorological influence (temperature,
        relative humidity) scatter plots to report
        """
        fig_name = self.sensor_name + '_normalized_' + self.eval_param \
            + '_met_report_fmt'

        # Search for figure created today
        try:
            fig_name += '_' + self.today + '.png'
            fig_path = self.figure_path + '\\'.join(
                                                (self.eval_param, fig_name))

            figure = open(fig_path, 'r')
            figure.close()

        except FileNotFoundError:
            self.plot_met_influence(report_fmt=True,
                                    plot_error_bars=False)

        metinf_loc = self.fig_locs['MetInfl']
        self.metinf = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(metinf_loc['left']),
                                    top=ppt.util.Inches(metinf_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metinf._element)

    def GetShape(self, slide_idx, shape_id=None, shape_loc=None):
        """
        Return shape object for tables based on known shape ID. Allows
        editing, modifying the table and its cells.

        Return either based on left and top location passed in inches to
        function (shape_loc=(left, top)), or by passing shape index to function
        """
        if shape_loc is not None:
            shp_l = shape_loc[0]
            shp_t = shape_loc[1]
            for shape in self.rpt.slides[slide_idx].shapes:
                if (math.isclose(shape.left.inches, shp_l, rel_tol=.05)
                   and math.isclose(shape.top.inches, shp_t, rel_tol=.05)):
                    return shape
            print('No shape within specified tolerance!')

        else:
            for shape in self.rpt.slides[slide_idx].shapes:
                if shape.shape_id == shape_id:
                    return shape

    def EditHeader(self):
        """
        Shape name                  Slide number         Shape ID
        ----------                  -----------     -------------------
        Report title                     1          35 (PM2.5),  9 (O3)
        Report title                     2          21 (PM2.5), 21 (O3)
        Report title                     3          13 (PM2.5), 15 (O3)
        Deployment, contact info         1          34 (PM2.5), 33 (O3)
        Deployment, contact info         2          20 (PM2.5), 22 (O3)
        Deployment, contact info         3          12 (PM2.5), 16 (O3)
        Photo placeholder                1           2 (PM2.5),  2 (O3)
        Photo placeholder                2           4 (PM2.5),  3 (O3)
        Photo placeholder                3          14 (PM2.5),  3 (O3)
        """
        # Title location for header
        title_left = 0.98
        title_top = 0.51
        text_vspace = 0
        # Get pptx text shape for modifying cells
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            title = self.GetShape(slide_idx, shape_loc=(title_left, title_top))

            # Set evaluation report title line 1: Parameter and test type
            title_line_1 = title.text_frame.paragraphs[0]
            run1 = title_line_1.add_run()
            run1.text = 'Testing Report - '
            run2 = title_line_1.add_run()
            run2.text = self.param_formatting_dict[
                                        self.eval_param]['baseline']
            run3 = title_line_1.add_run()
            run3.text = self.param_formatting_dict[
                                        self.eval_param]['subscript']
            font = run3.font
            font._element.set('baseline', '-25000')
            run4 = title_line_1.add_run()
            run4.text = ' Base Testing'

#            title_line_1.text = 'Testing Report - ' \
#                + self.eval_param + ' Base Testing'
            self.FormatText(title_line_1, alignment='left',
                            font_name='Calibri', font_size=30)
            title_line_1_font = title_line_1.font
            title_line_1_font.line_spacing = ppt.util.Pt(text_vspace)

            # Set evaluation report title line 2: Sensor name
            title_text_obj = title.text_frame.add_paragraph()
            title_text_obj.text = self.fmt_sensor_name
            self.FormatText(title_text_obj, alignment='left',
                            font_name='Calibri Light', font_size=30)
            title_font_obj = title_text_obj.font
            title_font_obj.line_spacing = ppt.util.Pt(text_vspace)

        # Header testing information location
        tester_left = 7.80
        tester_top = 0.36
        text_vspace = 10
        # Get pptx text shape for deployment and contact information
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            tester_info = self.GetShape(slide_idx, shape_loc=(tester_left,
                                                              tester_top))

            # Set deployment number
            tester_text = tester_info.text_frame.paragraphs[0]
            tester_text.text = self.testing_org['Deployment number']
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri', font_size=20, bold=True)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Set testing organization name (line 2)
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = self.testing_org['Org name'][0]
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Set testing organization name (line 2)
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = self.testing_org['Org name'][1]
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add contact email
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = (self.testing_org['Contact email'])
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add contact phone number
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = ('      ' + self.testing_org['Contact phone'])
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add current month to header
            month_year = dt.datetime.now().strftime('%B %Y')
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = month_year
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

        # Edit header photo (edits the image placeholder)
        pic_left = 12.56
        pic_top = 0.29
        # Get pptx text shape for deployment and contact information
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            pic = self.GetShape(slide_idx, shape_loc=(pic_left, pic_top))
            pic_path = self.figure_path + '\\deployment\\' + \
                self.sensor_name + '.png'
            if not os.path.exists(pic_path):
                sys.exit('No deployment picture found at', pic_path)
            else:
                pic.insert_picture(pic_path)

    def EditSiteTable(self):
        """
        Add details to testing organzation and site information table.

        Table name                  Table ID
        ----------                  --------
        Testing org, site info         18
        """
        # Get pptx table shape for modifying cells
        shape = self.GetShape(slide_idx=0, shape_id=18)

        # ------------- Cell 1: Testing organization information --------------
        cell = shape.table.cell(1, 1)

        # Add organization name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = (self.testing_org['Org name'][0] + ' - ' +
                         self.testing_org['Org name'][1])
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=14)

        # Add contact link
        text_obj = cell.text_frame.add_paragraph()
        run = text_obj.add_run()
        run.text = self.testing_org['Website']['website name']
        link = self.testing_org['Website']['website link']
        hlink = run.hyperlink
        hlink.address = link
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=11)

        # Add contact email, phone number
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = (self.testing_org['Contact email'] + '\n' + '      ' +
                         self.testing_org['Contact phone'])
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=11)

        # ----------- Cell 2: Testing location information ----------------
        cell = shape.table.cell(2, 1)

        # Add site name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.testing_loc['Site name']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # Add site address
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = self.testing_loc['Site address']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=11)

        # Add site lat long
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = (self.testing_loc['Site lat'] + ', ' +
                         self.testing_loc['Site long'])
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=11)

        # --------------- Cell 2: Testing location AQS ID ---------------------
        cell = shape.table.cell(3, 1)

        # Add site name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.testing_loc['Site AQS ID']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # ------------------ Cell 3: Sampling timeframe -----------------------
        cell = shape.table.cell(4, 1)

        # Add sampling timeframe for each group in deployment
        for i, grp in enumerate(self.tframe):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.tframe.keys())[i]
            grp_tframe = self.tframe[grp_name]
            text_obj.text = grp_name + ': ' + grp_tframe
            self.FormatText(text_obj, alignment='center', font_name='Calibri',
                            font_size=11)

    def EditSensorTable(self):
        """
        Add information to sensor information table

        Table name                  Table ID
        ----------                  --------
        Sensor info             49 (PM2.5), 30 (O3)
        """
        # Get pptx table shape for modifying cells
        if self.eval_param == 'PM25':
            shape = self.GetShape(slide_idx=0, shape_id=49)
        if self.eval_param == 'O3':
            shape = self.GetShape(slide_idx=0, shape_id=30)

        # Populate list with configured sensor recording interval(s)
        rec_interval = []
        for grp in self.deploy_dict['Deployment Groups']:
            grp_dict = self.deploy_dict['Deployment Groups'][grp]
            grp_sensors = grp_dict['sensors']
            for sensor in grp_sensors:
                interval = grp_sensors[sensor]['recording_interval']
                interval = interval.replace('.0', '')
                rec_interval.append(interval)
        rec_interval = list(set(rec_interval))
        rec_str = ', '.join(string for string in rec_interval)

        # ------------- Cell 1: Sensor manufacturer and model -----------------
        cell = shape.table.cell(1, 1)

        # Add sensor manufacturer, model name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.fmt_sensor_name
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # --------------- Cell 3: Sensor recording interval -------------------
        cell = shape.table.cell(3, 1)

        # Add recording interval
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = rec_str
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # -------------------- Sensor Serial ID cells -------------------------
        rows = [4, 5, 6]  # row indicies for sensor serial cells
        cols = [1, 3, 4]  # column indicies for sensor serial cells
        for i, iloc in enumerate(rows):
            for j, jloc in enumerate(cols):
                cell_n = 3*i + (j + 1)

                # Only fill in ID info for number of sensors in evaluation.
                # If number of sensors is, say, 8, then the last cell for the
                # 3x3 grid for serial numbers is left empty.
                if len(self.serial_grp_dict) < cell_n:
                    pass
                else:
                    cell = shape.table.cell(iloc, jloc)

                    serial_idx = cell_n - 1
                    sensor_serial = list(self.serial_grp_dict.keys()
                                         )[serial_idx]
                    sensor_grp = list(self.serial_grp_dict.values()
                                      )[serial_idx]

                    # Add group number
                    grp_obj = cell.text_frame.paragraphs[0]
                    grp_obj.text = sensor_grp
                    self.FormatText(grp_obj, alignment='center',
                                    font_name='Calibri', font_size=10.5)

                    # Add sensor serial ID
                    serial_obj = cell.text_frame.add_paragraph()
                    serial_obj.text = sensor_serial
                    self.FormatText(serial_obj, alignment='center',
                                    font_name='Calibri', font_size=10.5)

        # Check for issues with deployment
        for i, grp in enumerate(self.deploy_dict['Deployment Groups']):
            deploy_grp_data = (self.deploy_dict['Deployment Groups']
                                               [grp]['sensors'])
            grp_sensors = list(deploy_grp_data.keys())
            grp_status = [deploy_grp_data[str(j)]['deploy_issues']
                          for j in grp_sensors]

            cell = shape.table.cell(7, 2)
            if list(set(grp_status)) == ['False']:
                if i == 0:
                    textobj = cell.text_frame.paragraphs[0]
                else:
                    textobj = cell.text_frame.add_paragraph()
                textobj.text = grp + ': No Issues'
                self.FormatText(textobj, alignment='center',
                                font_name='Calibri', font_size=12)
            else:
                if i == 0:
                    textobj = cell.text_frame.paragraphs[0]
                else:
                    textobj = cell.text_frame.add_paragraph()
                textobj.text = grp + ': Issues with deployment'
                self.FormatText(textobj, alignment='center',
                                font_name='Calibri', font_size=12)

    def EditRefTable(self):
        """
        Add details to reference information table

        Table name                  Table ID
        ----------                  --------
        Reference info                 51
        """
        # Get pptx table shape for modifying cells
        shape = self.GetShape(slide_idx=0, shape_id=51)

        # ------------- Cell 1: Sensor manufacturer and model -----------------
        cell = shape.table.cell(1, 1)

        # Add reference manufacturer and model
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.ref_name
        self.FormatText(text_obj, alignment='center',
                        font_name='Calibri', font_size=14)

    def EditRefConcTable(self):
        """
        Add details to tables containing informtaion about reference
        concentrations. Located in different boxes based on the evaluation
        parameter type.

        Scatter plots box (PM2.5 only)
        Table name                  Table ID
        ----------                  --------
        Reference conc info          75 (PM25)

        Time series box (O3 only)
        ------------------------------------
        Table name                  Table ID
        ----------                  --------
        Reference conc info          56 (O3)
        """
        # Get pptx table shape for modifying cells
        if self.eval_param == 'PM25':
            shape = self.GetShape(slide_idx=0, shape_id=75)
        if self.eval_param == 'O3':
            shape = self.GetShape(slide_idx=0, shape_id=56)

        grp_info = self.deploy_dict['Deployment Groups']
        # reference concentration range
        self.refconc = {}
        for grp in list(grp_info.keys()):
            ref = 'Reference'
            try:
                self.refconc[grp] = \
                    '{0:3.1f}-{1:3.1f} (1-hr), '\
                    '{2:3.1f}-{3:3.1f} (24-hr)'.format(
                    grp_info[grp][self.eval_param][ref]['conc_min_hourly'],
                    grp_info[grp][self.eval_param][ref]['conc_max_hourly'],
                    grp_info[grp][self.eval_param][ref]['conc_min_daily'],
                    grp_info[grp][self.eval_param][ref]['conc_max_daily'])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # Number of periods reference exceeded concentration target
        if self.eval_param == 'PM25':
            exceed_str = 'n_exceed_conc_goal_daily'
        if self.eval_param == 'O3':
            exceed_str = 'n_exceed_conc_goal_hourly'

        self.refexceed = {}
        for grp in list(grp_info.keys()):
            try:
                self.refexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][self.eval_param]['Reference'][exceed_str])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # ------------- Cell 1: Range of ref concentrations -----------------
        cell = shape.table.cell(0, 1)
        for i, grp in enumerate(self.refconc):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.refconc.keys())[i]
            grp_refconc = self.refconc[grp_name]
            text_obj.text = grp_name + ': ' + grp_refconc
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=9)

        # ------------- Cell 2: N periods meeting conc. target-----------------
        cell = shape.table.cell(1, 1)
        for i, grp in enumerate(self.refexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.refexceed.keys())[i]
            grp_refexceed = self.refexceed[grp_name]
            text_obj.text = grp_name + ': ' + grp_refexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=9)

    def EditMetCondTable(self):
        """
        Met conditions box
        ---------------------------------------
        Table name                     Table ID
        ----------                     --------
        N outside target criteria   45 (O3), 74 (PM25)

        """
        # Get pptx table shape for modifying cells
        if self.eval_param == 'PM25':
            shape = self.GetShape(slide_idx=0, shape_id=74)
        if self.eval_param == 'O3':
            shape = self.GetShape(slide_idx=0, shape_id=45)

        grp_info = self.deploy_dict['Deployment Groups']

        # Number of 24-hr periods temp exceeded target criteria
        self.tempexceed = {}
        exceed_str = 'n_exceed_target_criteria_daily'
        met_conds = 'Meteorological Conditions'  # abbreviating
        temp = 'Temperature'  # abbreviating
        for grp in list(grp_info.keys()):
            try:
                self.tempexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][met_conds][temp][exceed_str])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # Number of 24-hr periods temp exceeded target criteria
        self.rhexceed = {}
        exceed_str = 'n_exceed_target_criteria_daily'
        met_conds = 'Meteorological Conditions'  # abbreviating
        rh = 'Relative Humidity'  # abbreviating
        for grp in list(grp_info.keys()):
            try:
                self.rhexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][met_conds][rh][exceed_str])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # ----------- Cell 1: N periods outside temp target criteria ----------
        cell = shape.table.cell(0, 1)
        for i, grp in enumerate(self.tempexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.tempexceed.keys())[i]
            grp_tempexceed = self.tempexceed[grp_name]
            text_obj.text = grp_name + ': ' + grp_tempexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=10)

        # ----------- Cell 2: N periods outside rh target criteria ------------
        cell = shape.table.cell(1, 1)
        for i, grp in enumerate(self.rhexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.rhexceed.keys())[i]
            grp_rhexceed = self.rhexceed[grp_name]
            text_obj.text = grp_name + ': ' + grp_rhexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=10)

    def EditMetInfTable(self):
        """
        Met influence box
        ------------------------------------
        Table name                  Table ID
        ----------                  --------
        N paired met conc vals   48 (O3), 76 (PM25)
        """
        # Get pptx table shape for modifying cells
        if self.eval_param == 'PM25':
            shape = self.GetShape(slide_idx=0, shape_id=76)
        if self.eval_param == 'O3':
            shape = self.GetShape(slide_idx=0, shape_id=48)

        grp_info = self.deploy_dict['Deployment Groups']

        # Number of 1-hr periods temp exceeded target criteria
        params = ['Temperature', 'Relative Humidity']
        met_conds = 'Meteorological Conditions'
        pair_str = 'n_measurement_pairs_hourly'

        dic = {}

        for i, param in enumerate(params):
            dic[param] = {grp: grp_info[grp][met_conds][param][pair_str]
                          for grp in list(grp_info.keys())}

            for j, grp in enumerate(self.deploy_dict['Deployment Groups']):
                cell = shape.table.cell(i, 1)
                if j == 0:
                    text_obj = cell.text_frame.paragraphs[0]
                else:
                    text_obj = cell.text_frame.add_paragraph()
                grp_val = dic[param][grp]
                text_obj.text = grp + ': ' + str(round(grp_val))
                self.FormatText(text_obj, alignment='center',
                                font_name='Calibri', font_size=10)

    def EditSensorRefTable(self, table):
        """
        """
        if self.eval_param == 'PM25':
            span_dict = {'Bias and Linearity': [1, 6],
                         'Data Quality': [7, 10],
                         'R^2': [12, 13],
                         'Slope': [14, 15],
                         'Intercept (b)\n(μg/m^3)': [16, 17],
                         'Uptime (%)': [18, 19],
                         'N pairs': [20, 21]}
            table_categories = {'1': 'Bias and Linearity',
                                '7': 'Data Quality'}
            metrics = {'12': 'R^2',
                       '14': 'Slope',
                       '16': 'Intercept\n(μg/m^3)',
                       '18': 'Uptime\n(%)',
                       '20': 'Number of paired\nsensor and '
                             'FRM/FEM\nconcentration pairs'}
            avg_intervals = {'23': '1-Hour',
                             '24': '24-Hour',
                             '25': '1-Hour',
                             '26': '24-Hour',
                             '27': '1-Hour',
                             '28': '24-Hour',
                             '29': '1-Hour',
                             '30': '24-Hour',
                             '31': '1-Hour',
                             '32': '24-Hour'}
            metric_targets = {'33': 'Metric Target Range',
                              '34': '≥ 0.70',
                              '35': '≥ 0.70',
                              '36': '1.0 ± 0.35',
                              '37': '1.0 ± 0.35',
                              '38': '-5 ≤ b ≤ 5',
                              '39': '-5 ≤ b ≤ 5',
                              '40': '75%*',
                              '41': '75%*',
                              '42': '-',
                              '43': '-'}

        if self.eval_param == 'O3':
            span_dict = {'Bias and Linearity': [1, 3],
                         'Data Quality': [4, 5]}
            table_categories = {'1': 'Bias and Linearity',
                                '4': 'Data Quality'}
            metrics = {'7': 'R^2',
                       '8': 'Slope',
                       '9': 'Intercept\n(ppbv)',
                       '10': 'Uptime\n(%)',
                       '11': 'Number of paired\nsensor and '
                             'reference\nconcentration pairs'}
            avg_intervals = {'13': '1-Hour',
                             '14': '1-Hour',
                             '15': '1-Hour',
                             '16': '1-Hour',
                             '17': '1-Hour'}
            metric_targets = {'18': 'Metric Target Range',
                              '19': '≥ 0.80',
                              '20': '1.0 ± 0.20',
                              '21': '-5 ≤ b ≤ 5',
                              '22': '75%*',
                              '23': '-'}

        self.grp_stats = self.stats_df.dropna()

        cells = self.SetSpanningCells(table, span_dict)

        if self.eval_param == 'O3':
            c1_row, c2_row, c3_row, c4_row, c5_row = 0, 0, 0, 0, 0
            h_stats = self.grp_stats.where(
                            self.grp_stats['Averaging Interval'] == 'Hourly'
                            ).dropna().reset_index(drop=True)

        if self.eval_param == 'PM25':
            (c1_row, c2_row, c3_row, c4_row, c5_row, c6_row, c7_row, c8_row,
             c9_row, c10_row) = (0, 0, 0, 0, 0, 0, 0, 0, 0, 0)
            h_stats = self.grp_stats.where(
                        self.grp_stats['Averaging Interval'] == 'Hourly'
                        ).dropna().reset_index(drop=True)
            d_stats = self.grp_stats.where(
                        self.grp_stats['Averaging Interval'] == 'Daily'
                        ).dropna().reset_index(drop=True)

        # Sensor group dictionary, reset group keys for indexing
        grp_df = (self.deploy_dict['Deployment Groups']
                                  [self.grp_name]['sensors'])
        grp_df = {i: grp_df[k] for i, k in enumerate(sorted(grp_df.keys()))}

        # Lists for storing metric values to compute mean across sensors
        rsqr_h_vals, rsqr_d_vals = [], []
        slope_h_vals, slope_d_vals = [], []
        intcpt_h_vals, intcpt_d_vals = [], []
        self.uptime_h_vals, self.uptime_d_vals = [], []
        n_h_vals, n_d_vals = [], []

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]

            n_header_rows = 4
            if self.eval_param == 'O3':
                datacols = 6
            if self.eval_param == 'PM25':
                datacols = 11
            datacell0 = datacols*n_header_rows
            datacellend = datacell0 + self.grp_n_sensors*datacols

            # Add table header info (rows 1-4)
            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if i == datacell0 + self.grp_n_sensors*datacols:
                # Row of avg values at bottom of table
                text_obj.text = 'Mean'

            # Add Sensor Serial IDs to column 1
            if (i % datacols == 0 and i >= datacell0
               and c1_row < self.grp_n_sensors):
                try:
                    text_obj.text = 'Sensor ' + grp_df[c1_row]['serial_id']
                except KeyError:
                    pass

            # Decision tree for adding sensor data to table
            if (i > datacell0 and i < datacell0 + (datacols)*self.grp_n_sensors
               and i % datacols != 0):
                j = i - (datacell0 + 1)

                # Metric column 1
                if j % datacols == 0:
                    try:
                        # 1-hr R^2
                        val = h_stats.loc[c1_row, 'R$^2$']
                        text_obj.text = format(val, '3.2f')
                        rsqr_h_vals.append(val)
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c1_row += 1

                # Metric column 2
                if (j-1) % datacols == 0:
                    try:
                        if self.eval_param == 'O3':
                            # 1-hr slope
                            val = h_stats.loc[c2_row, 'Slope']
                            slope_h_vals.append(val)
                        if self.eval_param == 'PM25':
                            # 24-hr R^2
                            val = d_stats.loc[c2_row, 'R$^2$']
                            rsqr_d_vals.append(val)

                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c2_row += 1

                # Metric column 3
                if (j-2) % datacols == 0:
                    try:
                        if self.eval_param == 'O3':
                            # 1-hr Intercept
                            val = h_stats.loc[c3_row, 'Intercept']
                            intcpt_h_vals.append(val)
                        if self.eval_param == 'PM25':
                            # 1-hr Slope
                            val = h_stats.loc[c3_row, 'Slope']
                            slope_h_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c3_row += 1

                # Metric column 4
                if (j-3) % datacols == 0:
                    # insert Uptime
                    try:
                        if self.eval_param == 'O3':
                            # 1-hr uptime
                            val = grp_df[c4_row]['uptime_hourly']
                            self.uptime_h_vals.append(val)
                        if self.eval_param == 'PM25':
                            # 24-hr Slope
                            val = d_stats.loc[c4_row, 'Slope']
                            slope_d_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c4_row += 1

                # Metric column 5
                if (j-4) % datacols == 0:
                    # insert N paired sensor/reference values
                    try:
                        if self.eval_param == 'O3':
                            # 1-hr N
                            val = h_stats.loc[c5_row, 'N']
                            n_h_vals.append(val)
                        if self.eval_param == 'PM25':
                            # 1-hr Intercept
                            val = h_stats.loc[c5_row, 'Intercept']
                            intcpt_h_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c5_row += 1

                # Columns 6-10 only in PM2.5 sensor-reference table
                if self.eval_param == 'PM25':

                    # Metric column 6
                    if (j-5) % datacols == 0:
                        try:
                            # 24-hr Intercept
                            val = d_stats.loc[c6_row, 'Intercept']
                            intcpt_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c6_row += 1

                    # Metric column 7
                    if (j-6) % datacols == 0:
                        try:
                            # 1-hr Uptime
                            val = grp_df[c7_row]['uptime_hourly']
                            self.uptime_h_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c7_row += 1

                    # Metric column 8
                    if (j-7) % datacols == 0:
                        try:
                            # 24-hr Uptime
                            val = grp_df[c8_row]['uptime_daily']
                            self.uptime_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c8_row += 1

                    # Metric column 9
                    if (j-8) % datacols == 0:
                        try:
                            # 1-hr N
                            val = h_stats.loc[c9_row, 'N']
                            n_h_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c9_row += 1

                    # Metric column 10
                    if (j-9) % datacols == 0:
                        try:
                            # 24-hr N
                            val = d_stats.loc[c10_row, 'N']
                            n_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c10_row += 1

            if i > datacellend:
                j = i - (datacellend + 1)

                # Metric Mean column 1
                if j % datacols == 0:
                    # Mean 1-hr R^2
                    if self.eval_param == 'O3':
                        val = np.mean(rsqr_h_vals)
                        txt = self.CheckTargets(rsqr_h_vals, metric='R$^2$')
                    # Mean 1-hr R^2
                    if self.eval_param == 'PM25':
                        val = np.mean(rsqr_h_vals)
                        txt = self.CheckTargets(rsqr_h_vals, metric='R$^2$')

                # Metric Mean column 2
                if (j - 1) % datacols == 0:
                    # Mean 1-hr Slope
                    if self.eval_param == 'O3':
                        val = np.mean(slope_h_vals)
                        txt = self.CheckTargets(slope_h_vals, metric='Slope')
                    # Mean 24-hr R^2
                    if self.eval_param == 'PM25':
                        val = np.mean(rsqr_d_vals)
                        txt = self.CheckTargets(rsqr_d_vals, metric='R$^2$')

                # Metric Mean column 3
                if (j - 2) % datacols == 0:
                    # Mean 1-hr Intercept
                    if self.eval_param == 'O3':
                        val = np.mean(intcpt_h_vals)
                        txt = self.CheckTargets(intcpt_h_vals,
                                                metric='Intercept')
                    # Mean 1-hr Slope
                    if self.eval_param == 'PM25':
                        val = np.mean(slope_h_vals)
                        txt = self.CheckTargets(slope_h_vals, metric='Slope')

                # Metric Mean column 4
                if (j - 3) % datacols == 0:
                    # Mean 1-hr Uptime
                    if self.eval_param == 'O3':
                        val = np.mean(self.uptime_h_vals)
                        txt = self.CheckTargets(self.uptime_h_vals,
                                                metric='Uptime')
                    # Mean 24-hr Slope
                    if self.eval_param == 'PM25':
                        val = np.mean(slope_d_vals)
                        txt = self.CheckTargets(slope_d_vals, metric='Slope')

                # Metric Mean column 5
                if (j - 4) % datacols == 0:
                    # Mean 1-hr N paired measurements
                    if self.eval_param == 'O3':
                        val = np.mean(n_h_vals)
                    # Mean 1-hr Intercept
                    if self.eval_param == 'PM25':
                        val = np.mean(intcpt_h_vals)
                        txt = self.CheckTargets(intcpt_h_vals,
                                                metric='Intercept')

                if self.eval_param == 'PM25':
                    # Metric Mean column 6
                    if (j - 5) % datacols == 0:
                        # Mean 24-hr Intercept
                        val = np.mean(intcpt_d_vals)
                        txt = self.CheckTargets(intcpt_d_vals,
                                                metric='Intercept')

                    # Metric Mean column 7
                    if (j - 6) % datacols == 0:
                        # Mean 1-hr Uptime
                        val = np.mean(self.uptime_h_vals)
                        txt = self.CheckTargets(self.uptime_h_vals,
                                                metric='Uptime')

                    # Metric Mean column 8
                    if (j - 7) % datacols == 0:
                        # Mean 24-hr Uptime
                        val = np.mean(self.uptime_d_vals)
                        txt = self.CheckTargets(self.uptime_d_vals,
                                                metric='Uptime')

                    # Metric Mean column 9
                    if (j - 8) % datacols == 0:
                        # Mean 1-hr N paired measurements
                        val = np.mean(n_h_vals)

                    # Metric Mean column 10
                    if (j - 9) % datacols == 0:
                        # Mean 24-hr N paired measurements
                        val = np.mean(n_d_vals)

                # Indicate number of sensors meeting performance metric target
                if txt is not None:
                    trgt_cell = cells[i - datacols*(self.grp_n_sensors + 2)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

                text_obj.text = format(val, '3.2f')

            # Configure text formatting
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)
            txt = None

    def EditErrorTable(self, table):

        error_stats = (self.deploy_dict['Deployment Groups']
                                       [self.grp_name]
                                       [self.eval_param]
                                       ['Error'])

        for key, value in error_stats.items():
            if value is None:
                error_stats[key] = -999

        if self.eval_param == 'PM25':
            datacols = 4
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Error': [1, 4],
                         'RMSE (μg/m^3)': [6, 7],
                         'nRMSE (%)': [8, 9]}
            table_categories = {'1': 'Error'}
            metrics = {'6': 'RMSE (μg/m^3)',
                       '8': 'nRMSE (%)'}
            avg_intervals = {'11': '1-Hour',
                             '12': '24-Hour',
                             '13': '1-Hour',
                             '14': '24-Hour'}
            metric_targets = {'15': 'Metric Target Range',
                              '16': '≤ 7',
                              '17': '≤ 7',
                              '18': '≤ 30',
                              '19': '≤ 30',
                              '20': 'Deployment Value'}
            metric_vals = {'21': format(error_stats['rmse_hourly'], '3.2f'),
                           '22': format(error_stats['rmse_daily'], '3.2f'),
                           '23': format(error_stats['nrmse_hourly'], '3.2f'),
                           '24': format(error_stats['nrmse_daily'], '3.2f')}
            for key, value in metric_vals.items():
                if value == '-999.00':
                    metric_vals[key] = ''

        if self.eval_param == 'O3':
            datacols = 2
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Error': [1, 2]}
            table_categories = {'1': 'Error'}
            metrics = {'4': 'RMSE (ppbv)',
                       '5': 'nRMSE (%)'}
            avg_intervals = {'7': '1-Hour',
                             '8': '1-Hour'}
            metric_targets = {'9': 'Metric Target Range',
                              '10': '≤ 7',
                              '11': '≤ 30',
                              '12': 'Deployment Value'}
            metric_vals = {'13': format(error_stats['rmse_hourly'], '3.2f'),
                           '14': format(error_stats['nrmse_hourly'], '3.2f')}
            for key, value in metric_vals.items():
                if value == '-999.00':
                    metric_vals[key] = ''

        cells = self.SetSpanningCells(table, span_dict)

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]

            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if str(i) in metric_vals:
                text_obj.text = metric_vals[str(i)]

            # Metric column 1
            if i > headercellend:
                if (i - 1) % datacols == 0:
                    # Mean 1-hr RMSE
                    if self.eval_param == 'O3':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')
                    # Mean 1-hr RMSE
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')
                # Metric column 2
                if (i - 1) % datacols == 1:
                    # 1-hr nRMSE
                    if self.eval_param == 'O3':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='nRMSE')
                    # 24-hr RMSE
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')

                if self.eval_param == 'PM25':
                    # Metric column 3
                    if (i - 1) % datacols == 2:
                        # 1-hr nRMSE
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='nRMSE')
                    # Metric column 4
                    if (i - 1) % datacols == 3:
                        # 24-hr nRMSE
                        if self.eval_param == 'PM25':
                            val = float(metric_vals[str(i)])
                            txt = self.CheckTargets(val, metric='nRMSE')

                # Indicate whether sensors meet performance metric target
                if txt is not None:
                    trgt_cell = cells[i - 2*(datacols + 1)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)

            txt = None

    def EditSensorSensorTable(self, table):

        grp_stats = (self.deploy_dict['Deployment Groups']
                                     [self.grp_name]
                                     [self.eval_param]
                                     ['Precision'])

        if self.eval_param == 'PM25':
            datacols = 8
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Precision (between collocated sensors)': [1, 4],
                         'Data Quality': [5, 8],
                         'CV\n(%)': [10, 11],
                         'SD\n(μg/m^3)': [12, 13],
                         'Uptime\n(%)': [14, 15],
                         'Number concurrent sensor values': [16, 17]}
            table_categories = {'1': 'Precision (between collocated sensors)',
                                '5': 'Data Quality'}
            metrics = {'10': 'CV\n(%)',
                       '12': 'Standard Deviation\n(μg/m^3)',
                       '14': 'Uptime\n(%)',
                       '16': 'Number of paired\nsensor and '
                             'reference\nconcentration pairs'}
            avg_intervals = {'19': '1-Hour',
                             '20': '24-Hour',
                             '21': '1-Hour',
                             '22': '24-Hour',
                             '23': '1-Hour',
                             '24': '24-Hour',
                             '25': '1-Hour',
                             '26': '24-Hour'}
            metric_targets = {'27': 'Metric Target Range',
                              '28': '≤ 30',
                              '29': '≤ 30',
                              '30': '≤ 5',
                              '31': '≤ 5',
                              '32': '75%*',
                              '33': '75%*',
                              '34': '-',
                              '35': '-',
                              '36': 'Deployment Value'}
            metric_vals = {'37': format(grp_stats['cv_hourly'], '3.2f'),
                           '38': format(grp_stats['cv_daily'], '3.2f'),
                           '39': format(grp_stats['std_hourly'], '3.2f'),
                           '40': format(grp_stats['std_daily'], '3.2f'),
                           '41': format(np.mean(self.uptime_h_vals), '3.2f'),
                           '42': format(np.mean(self.uptime_d_vals), '3.2f'),
                           '43': format(grp_stats['n_hourly']),
                           '44': format(grp_stats['n_daily'])}

        if self.eval_param == 'O3':
            datacols = 4
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Precision (between collocated sensors)': [1, 2],
                         'Data Quality': [3, 4]}

            table_categories = {'1': 'Precision (between collocated sensors)',
                                '3': 'Data Quality'}
            metrics = {'6': 'CV\n(%)',
                       '7': 'Standard Deviation\n(ppbv)',
                       '8': 'Uptime\n(%)',
                       '9': 'Number of paired\nsensor and '
                            'reference\nconcentration pairs'}
            avg_intervals = {'11': '1-Hour',
                             '12': '1-Hour',
                             '13': '1-Hour',
                             '14': '1-Hour'}
            metric_targets = {'15': 'Metric Target Range',
                              '16': '≤ 30',
                              '17': '≤ 5',
                              '18': '75%*',
                              '19': '-',
                              '20': 'Deployment Value'}
            metric_vals = {'21': format(grp_stats['cv_hourly']),
                           '22': format(grp_stats['std_hourly']),
                           '23': format(np.mean(self.uptime_h_vals), '3.2f'),
                           '24': format(grp_stats['n_hourly'])}

        cells = self.SetSpanningCells(table, span_dict)

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]
            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if str(i) in metric_vals:
                text_obj.text = metric_vals[str(i)]

            # Metric column 1
            if i > headercellend:
                j = i - headercellend
                if (j - 1) % datacols == 0:
                    # Mean 1-hr CV
                    if self.eval_param == 'O3':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                    # Mean 1-hr CV
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 2
                if (j - 1) % datacols == 1:
                    # 1-hr Std dev
                    if self.eval_param == 'O3':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='SD')
                    # 24-hr CV
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 3
                if (j - 1) % datacols == 2:
                    # 1-hr Uptime
                    if self.eval_param == 'O3':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')
                    # 1-hr Std dev
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 4
                if (j - 1) % datacols == 3:
                    # 24-hr Std dev
                    if self.eval_param == 'PM25':
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='SD')

                if self.eval_param == 'PM25':
                    # Metric column 5
                    if (j - 1) % datacols == 4:
                        # 1-hr uptime
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')
                    # Metric column 6
                    if (j - 1) % datacols == 5:
                        # 24-hr uptime
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')

                # Indicate whether sensors meet performance metric target
                if txt is not None:
                    trgt_cell = cells[i - 2*(datacols + 1)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)

            txt = None

    def SetSpanningCells(self, table, span_dict):

        cells = [cell for cell in table.iter_cells()]

        if span_dict is not None:
            for entry in span_dict:
                merge_start_idx = span_dict[entry][0]
                merge_end_idx = span_dict[entry][1]

                merge_start_cell = cells[merge_start_idx]
                merge_end_cell = cells[merge_end_idx]
                merge_start_cell.merge(merge_end_cell)

        return cells

    def EditTabularStats(self):
        """
        """
        self.n_grps = len(set(self.serial_grp_dict.values()))
        self.n_sensors = len(set(self.serial_grp_dict))

        # Use slide layout for generating additional slides
        tabular_layout_idx = 1
        tabular_layout = self.rpt.slide_layouts[tabular_layout_idx]

        legend_txt = {
                'line1': 'Single-valued metrics '
                         '(computed via entire evaluation dataset)',
                'line2': '○  Indicates that the metric value is not '
                         'within the target range',
                'line3': '●  Indicates that the metric value is within '
                         'the target range',
                'line4': '',
                'line5': 'Device-specific metrics (computed for each '
                         'sensor in evaluation)',
                'line6': '○○○	Metric value for none of devices tested '
                         'falls within the target range',
                'line7': '●○○	Metric value for one of devices tested '
                         'falls within the target range',
                'line8': '●●○   Metric value for two of devices tested '
                         'falls within the target range',
                'line9': '●●●   Metric value for three of devices tested '
                         'falls within the target range'
                }

        # List of unique testing groups
        grps = sorted(list(set(self.serial_grp_dict.values())))

        for grp_n, grp_name in enumerate(grps, 1):
            self.grp_n_sensors = list(
                        self.serial_grp_dict.values()).count(grp_name)
            self.grp_name = grp_name

            # Create new tabular stats page
            tabular_slide = self.rpt.slides.add_slide(
                                                tabular_layout)

            # Sensor Reference Table
            sr_frame, sr_table = self.CreateStatsTable(
                                        tabular_slide,
                                        table_type='sensor_reference')
            self.EditSensorRefTable(sr_table)

            # Error Table
            e_frame, e_table = self.CreateStatsTable(
                                        tabular_slide,
                                        table_type='error')
            self.EditErrorTable(e_table)

            # Intersensor (sensor-sensor) table
            ss_frame, ss_table = self.CreateStatsTable(
                                        tabular_slide,
                                        table_type='sensor_sensor')
            self.EditSensorSensorTable(ss_table)

            # Adjust sensor-reference table vertical position
            EMU_to_in = 1/914400.0
            sr_top = 4.03  # in
            sr_height = EMU_to_in*sr_frame.height
            sr_frame.top = ppt.util.Inches(sr_top)

            # Adjust error table vertical position
            e_t = 4.03 + sr_height + 0.4  # in
            e_h = EMU_to_in*e_frame.height
            e_frame.top = ppt.util.Inches(e_t)

            # Adjust Sensor-Sensor table vertical position
            ss_top = e_t + e_h + 2*0.4  # in
            ss_frame.top = ppt.util.Inches(ss_top)

            # Add section header
            tabular_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(2.81),  # top
                                ppt.util.Inches(3.16),  # width
                                ppt.util.Inches(0.47))  # height
            tabular_header_obj = tabular_header.text_frame.paragraphs[0]

            tabular_header_obj.text = 'Tabular Statistics'
            if self.n_grps > 1:
                tabular_header_obj.text += ' - ' + self.grp_name

            self.FormatText(tabular_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=22)

            # Add Sensor-Reference section header text label
            sr_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(3.30),  # top
                                ppt.util.Inches(12.81),  # width
                                ppt.util.Inches(0.44))  # height
            sr_header_obj = sr_header.text_frame.paragraphs[0]
            sr_header_obj.text = 'Sensor-Reference Correlation'
            self.FormatText(sr_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=20)

            # Add Sensor-Sensor section header text label
            ss_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(e_t + e_h + 0.12),  # top
                                ppt.util.Inches(12.81),  # width
                                ppt.util.Inches(0.44))  # height
            ss_header_obj = ss_header.text_frame.paragraphs[0]
            ss_header_obj.text = 'Sensor-Sensor Precision'
            self.FormatText(ss_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=20)

            connector1 = tabular_slide.shapes.add_connector(
                                ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                                ppt.util.Inches(0.85),  # start_x
                                ppt.util.Inches(3.30),  # start_y
                                ppt.util.Inches(16.24),  # end_x
                                ppt.util.Inches(3.30))  # end_y

            connector1.line.fill.solid()
            connector1.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171,
                                                                         171,
                                                                         171)

            connector2 = tabular_slide.shapes.add_connector(
                            ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                            ppt.util.Inches(0.85),  # start_x
                            ppt.util.Inches(e_t + e_h + 0.56),  # start_y
                            ppt.util.Inches(16.24),  # end_x
                            ppt.util.Inches(e_t + e_h + 0.56))  # end_y

            connector2.line.fill.solid()
            connector2.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171,
                                                                         171,
                                                                         171)

            # Add Sensor-Sensor text label
            legend_h = 2.7
            legend = tabular_slide.shapes.add_shape(
                                ppt.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE,
                                ppt.util.Inches(8.83),  # left
                                ppt.util.Inches(e_t + 0.5*(e_h + - legend_h)),
                                ppt.util.Inches(6.69),  # width
                                ppt.util.Inches(legend_h))

            legend.fill.solid()
            legend.fill.fore_color.rgb = ppt.dml.color.RGBColor(236,
                                                                236,
                                                                240)

            legend.line.fill.solid()
            legend.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(214,
                                                                     216,
                                                                     226)
            legend.line.width = ppt.util.Pt(2.5)

            legend_obj = legend.text_frame.paragraphs[0]
            legend_obj.text = legend_txt['line1']
            self.FormatText(legend_obj, alignment='left',
                            font_name='Calibri', font_size=16, bold=False)

            for i, line in enumerate(legend_txt):
                if i > 0:
                    legend_obj = legend.text_frame.add_paragraph()
                    legend_obj.text = legend_txt[line]

                font = legend_obj.font
                font.name = 'Calibri'
                font.Bold = False
                font.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)

                if line == 'line4':
                    font.size = ppt.util.Pt(8)
                else:
                    font.size = ppt.util.Pt(15)

    def CreateStatsTable(self, slide, table_type='sensor_reference'):

        cell_margin = 0.001
        table_spacing = 0.4
        shapes = slide.shapes

        # Create Sensor-Reference Correlation Table ---------------------------
        if table_type == 'sensor_reference':
            if self.eval_param == 'PM25':
                nrows = 5 + self.grp_n_sensors
                ncols = 11
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(14.4)
                height = ppt.util.Inches(3.07 + 0.45*self.grp_n_sensors)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69)
                # grey out cells at index (from l-r where top left is 0)
                greyed_cells = [42, 43]

            if self.eval_param == 'O3':
                nrows = 5 + self.grp_n_sensors
                ncols = 6
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(14.4)
                height = ppt.util.Inches(3.07 + 0.45*self.grp_n_sensors)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69)
                greyed_cells = [23]

        if table_type == 'error':
            if self.eval_param == 'PM25':
                nrows = 5
                ncols = 5
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(7.2)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + table_spacing)
                greyed_cells = None

            if self.eval_param == 'O3':
                nrows = 5
                ncols = 3
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(7.2)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + table_spacing)
                greyed_cells = None

        if table_type == 'sensor_sensor':
            if self.eval_param == 'PM25':
                nrows = 5
                ncols = 9
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(12.0)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + 2*table_spacing + 3.23)
                greyed_cells = [34, 35]

            if self.eval_param == 'O3':
                nrows = 5
                ncols = 5
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(12.0)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + 2*table_spacing + 3.23)
                greyed_cells = [19]

        frame = shapes.add_table(nrows, ncols, left, top,
                                 width, height)
        table = frame.table
        table.horz_banding = False

        for cell_idx, cell in enumerate(table.iter_cells()):
            # Set cell border width
            self.SetCellBorder(cell)

            # Set dark blue cell fill color
            if cell_idx < 3*ncols + 1 or cell_idx % ncols == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ppt.dml.color.RGBColor(191,
                                                                  208,
                                                                  235)
            # Set transparent (background) fill for top 3 cells
            if cell_idx % ncols == 0 and cell_idx <= 2*ncols:
                cell.fill.background()

            # Grey out cells where no target range specified
            if greyed_cells is not None:
                if cell_idx in greyed_cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ppt.dml.color.RGBColor(214,
                                                                      216,
                                                                      226)
            # Set cell margins to near zero
            cell.margin_left = ppt.util.Inches(cell_margin)
            cell.margin_right = ppt.util.Inches(cell_margin)
            cell.margin_top = ppt.util.Inches(cell_margin)
            cell.margin_bottom = ppt.util.Inches(cell_margin)
            # Vertical position text in middle of cells
            cell.vertical_anchor = ppt.enum.text.MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = ppt.util.Pt(14)

        # Modify table column width
        table.columns[0].width = ppt.util.Inches(2.2)
        for col_idx in np.arange(1, len(table.columns)):
            table.columns[col_idx].width = col_width

        # Modify table header row height
        table.rows[0].height = ppt.util.Inches(0.47)  # Category
        table.rows[1].height = ppt.util.Inches(0.97)  # Metric label
        table.rows[2].height = ppt.util.Inches(0.72)  # Averaging interval
        for row_idx in np.arange(3, len(table.rows)):
            table.rows[row_idx].height = ppt.util.Inches(0.45)

        return frame, table

    def PrintPpptxShapes(self, slide_number=1, shape_type='all'):
        """
        Diagnostic tool for indicating shape ids and locations on reporting
        template slides
        """
        print("{:^6s}{:^18s}{:^12s}{:^10s}".format('ID', 'Type',
                                                   'Left loc', 'Top loc'))
        print("{:^6s}{:^18s}{:^12s}{:^10s}".format('', '', '(in)', '(in)'))
        print('{:^45s}'.format(45*'-'))
        if shape_type == 'all':
            for shape in self.rpt.slides[slide_number - 1].shapes:
                print("%4s %16s %10.2f %10.2f" %
                      (shape.shape_id, shape.shape_type,
                       shape.left.inches, shape.top.inches))
        else:
            for shape in self.rpt.slides[slide_number - 1].shapes:
                if str(shape.shape_type).startswith(shape_type):
                    print("%4s %16s %10.2f %10.2f" %
                          (shape.shape_id, shape.shape_type,
                           shape.left.inches, shape.top.inches))

    def SubElement(self, parent, tagname, **kwargs):
        """
        Code for editing tabular cell border width. Based on Steve Canny's code
        at the following link:
            https://groups.google.com/g/python-pptx/c/UTkdemIZICw
        """
        element = ppt.oxml.xmlchemy.OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def SetCellBorder(self, cell, border_color="ffffff", border_width='20000'):
        """
        Code for editing tabular cell border width. Based on Steve Canny's code
        at the following links:
            https://groups.google.com/g/python-pptx/c/UTkdemIZICw
            https://stackoverflow.com/questions/42610829/
            python-pptx-changing-table-style-or-adding-borders-to-cells
        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = self.SubElement(tcPr, lines, w=border_width,
                                 cap='flat', mpd='sng', algn='ctr')
            solidFill = self.SubElement(ln, 'a:solidFill')
            srgbClr = self.SubElement(solidFill, 'a:srgbClr',
                                      val=border_color)
            prstDash = self.SubElement(ln, 'a:prstDash',
                                       val='solid')
            round_ = self.SubElement(ln, 'a:round')
            headEnd = self.SubElement(ln, 'a:headEnd',
                                      type='none', w='med', len='med')
            tailEnd = self.SubElement(ln, 'a:tailEnd',
                                      type='none', w='med', len='med')

    def SetSubscript(font):
        """
        Workaround for making font object text subscript (not included in
        python-pptx as of v0.6.19)

        Code via Martin Packer:
            https://stackoverflow.com/questions/61329224/
            how-do-i-add-superscript-subscript-text-to-powerpoint-using
            -python-pptx
        """
        font._element.set('baseline', '-25000')

    def SetSuperscript(font):
        """
        Workaround for making font object text superscript (not included in
        python-pptx as of v0.6.19)

        Code via Martin Packer:
            https://stackoverflow.com/questions/61329224/
            how-do-i-add-superscript-subscript-text-to-powerpoint-using
            -python-pptx
        """
        font._element.set('baseline', '30000')

    def MoveSlide(slides, slide, new_idx):
        """
        Move a slide from one index to another.

        Code via GitHub user Amazinzay (Feb 17 2021):
            https://github.com/scanny/python-pptx/issues/68
        """
        slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])


    def FormatText(self, text_obj, alignment='center', font_name='Calibri',
                   font_size=24, bold=False, italic=False):

        if alignment == 'center':
            text_obj.alignment = ppt.enum.text.PP_ALIGN.CENTER
        if alignment == 'left':
            text_obj.alignment = ppt.enum.text.PP_ALIGN.LEFT

        font_obj = text_obj.font
        font_obj.name = font_name
        font_obj.size = ppt.util.Pt(font_size)
        font_obj.bold = bold
        font_obj.italic = italic

    def CheckTargets(self, metric_vals, metric):
        # Place float / int values (single-valued intersensor stats) into
        # list for parsing
        if type(metric_vals) != list:
            metric_vals = [metric_vals]

        # Performance target values for PM2.5 base test evaluations
        if self.eval_param == 'PM25':
            Rsqr_min, Rsqr_max = 0.70, 1.0
            Slope_min, Slope_max = 0.65,  1.35
            Intcpt_min, Intcpt_max = -5, 5
            Uptime_min, Uptime_max = 90, 100
            RMSE_min, RMSE_max = 0, 7
            nRMSE_min, nRMSE_max = 0, 30
            CV_min, CV_max = 0, 30
            SD_min, SD_max = 0, 5

        # Performance target values for O3 base test evaluations
        if self.eval_param == 'O3':
            Rsqr_min, Rsqr_max = 0.80, 1.0
            Slope_min, Slope_max = 0.80,  1.20
            Intcpt_min, Intcpt_max = -5, 5
            Uptime_min, Uptime_max = 75, 100
            RMSE_min, RMSE_max = 0, 5
            nRMSE_min, nRMSE_max = 0, 30
            CV_min, CV_max = 0, 30
            SD_min, SD_max = 0, 5

        if metric == 'R$^2$':
            metric_min, metric_max = Rsqr_min, Rsqr_max
        if metric == 'Slope':
            metric_min, metric_max = Slope_min, Slope_max
        if metric == 'Intercept':
            metric_min, metric_max = Intcpt_min, Intcpt_max
        if metric == 'Uptime':
            metric_min, metric_max = Uptime_min, Uptime_max
        if metric == 'RMSE':
            metric_min, metric_max = RMSE_min, RMSE_max
        if metric == 'nRMSE':
            metric_min, metric_max = nRMSE_min, nRMSE_max
        if metric == 'CV':
            metric_min, metric_max = CV_min, CV_max
        if metric == 'SD':
            metric_min, metric_max = SD_min, SD_max

        # Number of sensors meeting target
        n_sensors = len(metric_vals)
        n_passed = sum(metric_min <= val <= metric_max for val
                       in metric_vals)

        open_dot = '○'
        closed_dot = '●'

        if n_sensors != 0:
            pcnt_passed = 100*(n_passed / n_sensors)

        if metric in ['R$^2$', 'Slope', 'Intercept', 'Uptime']:
            text = n_passed*closed_dot + (n_sensors - n_passed)*open_dot
            if n_sensors == 1 and metric == 'Uptime':
                if pcnt_passed == 100:
                    text = closed_dot
                else:
                    text = open_dot

        if metric in ['CV', 'SD', 'RMSE', 'nRMSE']:
            if pcnt_passed == 100:
                text = closed_dot
            elif n_sensors > 0:
                text = open_dot
            else:
                # Metric values empty?
                text = ''

        return text

    def CreateReport(self):
        """
        Existing figures are assumed to have been created on the same day of
        class instantiation. If a figure filename is not found, sensor data
        are loaded via the SensorEvaluation class and the figure is generated.

        Note that Edit Header should be the last routine, loops over slides to
        add header information.
        """
        print('Creating Testing Report for', self.sensor_name)

        # Set figure positions
        self.FigPositions()

        print('..Adding figures to report')
        # Add figures to report
        self.AddSingleScatterPlot()
        self.AddTimeseriesPlot()
        self.AddMetricsPlot()
        self.AddMetDistPlot()
        self.AddMetInflPlot()

        print('..Adding tabular data')
        # Modify report tables
        self.EditSiteTable()
        self.EditSensorTable()
        self.EditRefTable()
        self.EditRefConcTable()
        self.EditMetCondTable()
        self.EditMetInfTable()
        self.EditTabularStats()

        self.AddMultiScatter()
        self.EditHeader()

        # Move the supplemnental info table to the last slide position
        # Code via github user Amazinzay (Feb 17 2021):
        # https://github.com/scanny/python-pptx/issues/68
        slides = self.rpt.slides
        slide = slides[1]
        new_idx = len(slides)
        slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])

        self.SaveReport()

    def SaveReport(self):
        print('..Saving report')
        self.rpt_name = 'Base_Testing_Report_' + self.eval_param\
                        + '_' + self.sensor_name + '_' + self.today + '.pptx'

        save_dir = '\\'.join((self.lib_path, 'Reports',
                              self.sensor_name, self.eval_param))
        save_path = '\\'.join((save_dir, self.rpt_name))

        if not os.path.exists(save_dir):
            os.makedirs(save_dir)
            print('..Creating directory:')
            print('....' + save_dir)

        print('....' + save_path.replace(self.lib_path, ''))
        self.rpt.save(save_path)
