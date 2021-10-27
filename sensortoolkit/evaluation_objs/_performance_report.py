# -*- coding: utf-8 -*-
"""
Description.

================================================================================

@Author:
  | Samuel Frederick, NSSC Contractor (ORAU)
  | U.S. EPA / ORD / CEMM / AMCD / SFSB

Created:
  Tue Dec 15 08:53:19 2020
Last Updated:
  Mon Jun 28 16:17:59 2021
"""
import pptx as ppt
import datetime as dt
import numpy as np
import math
import os
import sys
import warnings
from sensortoolkit.evaluation_objs import SensorEvaluation


class PerformanceReport(SensorEvaluation):
    """Generates air sensor performance evaluation reports.

    Reports are intended for evaluations following U.S. EPA's
    recommendations for base testing of air sensors at outdoor ambient air
    monitoring sites and collocated alongside FRM/FEM monitors for use in NSIM
    applications.

    In February 2021, U.S. EPA released two reports detailing recommended
    performance testing protocols, metrics, and target values for the
    evaluation of sensors measuring either fine particulate matter (PM2.5)
    or ozone (O3). More detail about EPA's sensor evaluation research as well
    as both reports can be found online at EPA's Air Sensor Toolbox:
        https://www.epa.gov/air-sensor-toolbox

    NOTE -
    PerformanceReport is an inherited class of SensorEvaluation. As a result,
    it inherits all the class and instance attributes of SensorEvaluation,
    including its numerous variables and data structures. Programmatically,
    PerformanceReport is intended as a direct extension of SensorEvaluation;
    users can easily interact with all the attributes and data stuctures for
    sensor evaluations. However, whereas SensorEvaluation allows analysis of a
    wide number of pollutants and parameters, PerformanceReport is intended for
    constructing reports pertaining to sensors measuring either fine
    particulate matter (PM2.5) or ozone (O3) following U.S. EPA's
    recommended protocols and testing metrics for evaluating these sensors.
    Module will exit execution if parameters other than 'PM25' or 'O3' are
    specified.

    Args:
        sensor (AirSensor object):
            Description.
        param (Parameter object):
            Description.
        reference (ReferenceMethod object):
            Description.
        write_to_file (TYPE, optional): DESCRIPTION. Defaults to False.
        figure_search (TYPE, optional): DESCRIPTION. Defaults to False.
        **kwargs (TYPE): DESCRIPTION.

    Returns:
            None.

    """

    # Evaluation parameters for which the PerformanceReport class can
    # constuct reports
    report_params = ['PM25', 'O3']

    def __init__(self, sensor, param, reference=None, write_to_file=False,
                 figure_search=False, **kwargs):

        # Add keyword arguments (testing_loc, testing_org, etc.)
        self.__dict__.update(**kwargs)
        self.kwargs = kwargs

        # Inherit the SensorEvaluation class instance attributes
        super().__init__(sensor, param, reference, write_to_file,
                         **kwargs)

        if self._param_name not in self.report_params:
            sys.exit('Reporting template not configured for '
                     + self._param_name)

        self.figure_search = figure_search

        # Placeholder method for formatted sensor name, replace '_' with spaces
        self.fmt_sensor_name = self.kwargs.get('fmt_sensor_name',
                                               self.name.replace('_',
                                                                        ' '))

        self.today = dt.datetime.now().strftime('%y%m%d')

        self.template_name = ('Reporting_Template_Base_' + self._param_name
                              + '.pptx')
        # Path to reporting template
        self.template_path = os.path.abspath(os.path.join(__file__,
                        '../templates', self._param_name, self.template_name))

        # Details about testing and deployment site
        self.testing_org = self.kwargs.get('testing_org',
                                           self.deploy_dict[
                                                   'Testing Organization'])
        self.testing_loc = self.kwargs.get('testing_loc',
                                           self.deploy_dict[
                                                   'Testing Location'])

        # Populate deployment dictionary with performance metric results
        self.calculate_metrics()

        # Sampling timeframe
        self.tframe = {grp:
                       self.deploy_dict['Deployment Groups'][grp]['eval_start']
                       + ' - ' +
                       self.deploy_dict['Deployment Groups'][grp]['eval_end']
                       for grp in list(
                               self.deploy_dict['Deployment Groups'].keys())
                       }

        # Keys are sensor serial IDs, values are deployment group number
        # Useful if multiple evaluation groups deployed
        self.serial_grp_dict = {}
        for grp in self.deploy_dict['Deployment Groups']:
            grp_dict = self.deploy_dict['Deployment Groups'][grp]
            grp_sensors = grp_dict['sensors']
            for sensor in grp_sensors:
                serial = grp_sensors[sensor]['serial_id']
                self.serial_grp_dict[serial] = grp

        # Initialize report object
        self.rpt = ppt.Presentation(self.template_path)
        self.shapes = self.rpt.slides[0].shapes

        # Shape at backgroud around which to orient other figures
        self.cursor_sp = self.shapes[0]._element

        # The number of unique averaging intervals at which data will be
        # presented. Either param.averaging == ['1-hour'] (gases,
        # n_avg_intervals == 1) or param.averaging == ['1-hour', '24-hour']
        # (PM, n_avg_intervals == 2).
        self.n_avg_intervals = len(self.param.averaging)

        # Initialize figure positions in report
        self.FigPositions()

        # Plotting: determine the max concentration for average of concurrent
        # sensor measurements and also the ref max concentration. Select the
        # upper limit for plots as 1.25x the larger of these values.
        sensor_avg_cmax = self.avg_hrly_df['mean_' + self._param_name + '_Value'].max()
        ref_cmax = self.hourly_ref_df[self._param_name + '_Value'].max()
        self.plot_cmax = 1.25*max(sensor_avg_cmax, ref_cmax)

    def FigPositions(self):
        """Assign figure positions for reports.

        Values are in inches, specifying the left and top center location of
        each figure.

        Returns:
            None.

        """
        self.fig_locs = {'SingleScatter': {'left': '',
                                           'top': ''},
                         'TripleScatter': {'left': '',
                                           'top': ''},
                         'Timeseries': {'left': '',
                                        'top': ''},
                         'MetricPlot': {'left': '',
                                        'top': ''},
                         'MetDist': {'left': '',
                                     'top': ''},
                         'MetInfl': {'left': '',
                                     'top': ''}
                         }

        if self.n_avg_intervals == 2:
            self.fig_locs['SingleScatter']['left'] = 11.11
            self.fig_locs['SingleScatter']['top'] = 8.16
            self.fig_locs['TripleScatter']['left'] = 2.35
            self.fig_locs['TripleScatter']['top'] = 3.82
            self.fig_locs['Timeseries']['left'] = 0.63
            self.fig_locs['Timeseries']['top'] = 8.15
            self.fig_locs['MetricPlot']['left'] = 0.65
            self.fig_locs['MetricPlot']['top'] = 13.19
            self.fig_locs['MetDist']['left'] = 1.91
            self.fig_locs['MetDist']['top'] = 17.58
            self.fig_locs['MetInfl']['left'] = 8.24
            self.fig_locs['MetInfl']['top'] = 17.54

        if self.n_avg_intervals == 1:
            self.fig_locs['SingleScatter']['left'] = 11.57
            self.fig_locs['SingleScatter']['top'] = 8.14
            self.fig_locs['TripleScatter']['left'] = 2.35
            self.fig_locs['TripleScatter']['top'] = 3.82
            self.fig_locs['Timeseries']['left'] = 0.63
            self.fig_locs['Timeseries']['top'] = 8.19
            self.fig_locs['MetricPlot']['left'] = 0.7
            self.fig_locs['MetricPlot']['top'] = 12.78
            self.fig_locs['MetDist']['left'] = 1.6
            self.fig_locs['MetDist']['top'] = 17.3
            self.fig_locs['MetInfl']['left'] = 8.28
            self.fig_locs['MetInfl']['top'] = 17.31

    def FigureSearch(self, figure_name, subfolder=None):
        """Indicate whether a figure exists and the full path to the figure.

        Args:
            figure_name (TYPE): DESCRIPTION.
            subfolder (TYPE, optional): DESCRIPTION. Defaults to None.

        Returns:
            TYPE: DESCRIPTION.
            full_figure_path (TYPE): DESCRIPTION.

        """
        if subfolder is None:
            subfolder = self._param_name
        # Search for figure created today
        figure_name += '_' + self.today + '.png'
        full_figure_path = self.figure_path + '\\'.join((subfolder,
                                                         figure_name))

        return os.path.exists(full_figure_path), full_figure_path

    def AddSingleScatterPlot(self, **kwargs):
        """Add sensor vs. reference scatter plots to report.

        Args:
            **kwargs (dict): Keyword arguments passed to plot_sensor_scatter()
            subroutine for drawing scatter plots.

        Returns:
            None.

        """
        fig_name = self.name + '_vs_' + self.ref_name + '_report_fmt'

        fig_exists, fig_path = self.FigureSearch(fig_name)

        # Draw figure if no figure exists at path or if figure_search attrib
        # is false
        create_figure = False
        if self.figure_search is False:
            create_figure = True
        else:
            if fig_exists is False:
                create_figure = True

        if create_figure:
            self.plot_sensor_scatter(
                plot_subset=kwargs.get('plot_subset', ['1']),
                # plot_limits=kwargs.get('plot_limits', (-1, self.plot_cmax)),
                # tick_spacing=kwargs.get('tick_spacing', 5),
                # text_pos=kwargs.get('text_pos', 'upper_left'),
                report_fmt=True)

        scatter_loc = self.fig_locs['SingleScatter']
        self.scatterplt = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(scatter_loc['left']),
                                    top=ppt.util.Inches(scatter_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.scatterplt._element)

    def AddMultiScatter(self, **kwargs):
        """Add Sensor vs. reference scatter plots for all sensors.

        Args:
            **kwargs (dict): Keyword arguments passed to plot_sensor_scatter()
            subroutine for drawing scatter plots.

        Returns:
            None.

        """
        # Use slide layout for generating additional slides
        slide_layout_idx = 0
        slide_layout = self.rpt.slide_layouts[slide_layout_idx]
        # Create new page
        slide = self.rpt.slides.add_slide(slide_layout)

        # Add Sensor-Sensor section header text label
        section_header = slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(2.81),  # top
                                ppt.util.Inches(3.16),  # width
                                ppt.util.Inches(0.47))  # height
        section_header_obj = section_header.text_frame.paragraphs[0]
        section_header_obj.text = 'Sensor-FRM/FEM Scatter Plots'
        self.FormatText(section_header_obj, alignment='left',
                        font_name='Calibri Light', font_size=22)

        # A horizontal line separating the header text from the section body
        hline = slide.shapes.add_connector(
                            ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                            ppt.util.Inches(0.85),  # start_x
                            ppt.util.Inches(3.30),  # start_y
                            ppt.util.Inches(16.24),  # end_x
                            ppt.util.Inches(3.30))  # end_y

        hline.line.fill.solid()
        hline.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171, 171, 171)

        # Loop over averaging intervals specified for the parameter
        for i, averaging_interval in enumerate(self.param.averaging):
            if self.n_sensors > 1:
                plural = 's'
            else:
                plural = ''
            fig_name = (self.name + '_vs_' + self.ref_name +
                        '_' + averaging_interval + '_' + str(self.n_sensors) +
                        '_' + 'sensor' + plural)

            fig_exists, fig_path = self.FigureSearch(fig_name)

            # Draw figure if no figure exists at path or if figure_search
            # attrib is false
            create_figure = False
            if self.figure_search is False:
                create_figure = True
            else:
                if fig_exists is False:
                    create_figure = True

            if create_figure:
                self.plot_sensor_scatter(
                    averaging_interval,
                    plot_limits=kwargs.get('plot_limits',
                                           (-1, self.plot_cmax)),
                    # tick_spacing=kwargs.get('tick_spacing', 5),
                    # text_pos=kwargs.get('text_pos', 'upper_left')
                    )

            if self.n_sensors <= 3:
                scatter_loc = self.fig_locs['TripleScatter']
                fig_height = 5.62  # height of triple scatter figure in inches
            else:
                # Add a page, place figure on new page
                # FIXME: Correct figure position needed for scatter plots with
                # more than three sensors.
                scatter_loc = self.fig_locs['TripleScatter']

                # TODO: set correct figure height for figures with mult. rows
                fig_height = 5.62  # height of triple scatter figure in inches

            left = ppt.util.Inches(scatter_loc['left'])
            top = ppt.util.Inches(scatter_loc['top'] + i*fig_height)

            slide.shapes.add_picture(fig_path, left, top)

    def AddTimeseriesPlot(self, **kwargs):
        """Add timeseries plots to report.

        Args:
            **kwargs (TYPE): DESCRIPTION.

        Returns:
            None.

        """
        fig_name = self.name + '_timeseries_' + self._param_name \
            + '_report_fmt'

        fig_exists, fig_path = self.FigureSearch(fig_name)

        # Draw figure if no figure exists at path or if figure_search
        # attrib is false
        create_figure = False
        if self.figure_search is False:
            create_figure = True
        else:
            if fig_exists is False:
                create_figure = True

        if create_figure:
            self.plot_timeseries(report_fmt=True)

        timeseries_loc = self.fig_locs['Timeseries']
        self.timeseries = self.shapes.add_picture(
                                fig_path,
                                left=ppt.util.Inches(timeseries_loc['left']),
                                top=ppt.util.Inches(timeseries_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.timeseries._element)

    def AddMetricsPlot(self, **kwargs):
        """Add Performance target metric boxplots/dot plots to report.

        Args:
            **kwargs (TYPE): DESCRIPTION.

        Returns:
            None.

        """
        fig_name = self.name + '_regression_boxplot_' + self._param_name

        fig_exists, fig_path = self.FigureSearch(fig_name)

        # Draw figure if no figure exists at path or if figure_search
        # attrib is false
        create_figure = False
        if self.figure_search is False:
            create_figure = True
        else:
            if fig_exists is False:
                create_figure = True

        if create_figure:
            self.plot_metrics()

        metricplt_loc = self.fig_locs['MetricPlot']
        self.metricplot = self.shapes.add_picture(
                                fig_path,
                                left=ppt.util.Inches(metricplt_loc['left']),
                                top=ppt.util.Inches(metricplt_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metricplot._element)

    def AddMetDistPlot(self, **kwargs):
        """Add meteorological distribution (Temp, RH) to report.

        Args:
            **kwargs (TYPE): DESCRIPTION.

        Returns:
            None.

        """
        fig_name = self.name + '_met_distplot_report_fmt'

        fig_exists, fig_path = self.FigureSearch(fig_name, subfolder='Met')

        # Draw figure if no figure exists at path or if figure_search
        # attrib is false
        create_figure = False
        if self.figure_search is False:
            create_figure = True
        else:
            if fig_exists is False:
                create_figure = True

        if create_figure:
            self.plot_met_dist()

        metdist_loc = self.fig_locs['MetDist']
        self.metdist = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(metdist_loc['left']),
                                    top=ppt.util.Inches(metdist_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metdist._element)

    def AddMetInflPlot(self, **kwargs):
        """Add normalized met. influence scatter (Temp, RH) to report.

        Args:
            **kwargs (TYPE): DESCRIPTION.

        Returns:
            None.

        """
        fig_name = self.name + '_normalized_' + self._param_name \
            + '_met_report_fmt'

        fig_exists, fig_path = self.FigureSearch(fig_name)

        # Draw figure if no figure exists at path or if figure_search
        # attrib is false
        create_figure = False
        if self.figure_search is False:
            create_figure = True
        else:
            if fig_exists is False:
                create_figure = True

        if create_figure:
            self.plot_met_influence(report_fmt=True,
                                    plot_error_bars=False)

        metinf_loc = self.fig_locs['MetInfl']
        self.metinf = self.shapes.add_picture(
                                    fig_path,
                                    left=ppt.util.Inches(metinf_loc['left']),
                                    top=ppt.util.Inches(metinf_loc['top']))

        # Move image to 0 z-order (background)
        self.cursor_sp.addprevious(self.metinf._element)

    def GetShape(self, slide_idx, shape_id=None, shape_loc=None):
        """Retrieve shape object for tables based on known shape ID.

        Allows for editing, modifying the table and its cells.

        Return either based on left and top location passed in inches to
        function (shape_loc=(left, top)), or by passing shape index to function

        Args:
            slide_idx (TYPE): DESCRIPTION.
            shape_id (TYPE, optional): DESCRIPTION. Defaults to None.
            shape_loc (TYPE, optional): DESCRIPTION. Defaults to None.

        Returns:
            shape (TYPE): DESCRIPTION.

        """
        if shape_loc is not None:
            shp_l = shape_loc[0]
            shp_t = shape_loc[1]
            for shape in self.rpt.slides[slide_idx].shapes:
                if (math.isclose(shape.left.inches, shp_l, rel_tol=.05)
                   and math.isclose(shape.top.inches, shp_t, rel_tol=.05)):
                    return shape
            print('No shape within specified tolerance!')

        else:
            for shape in self.rpt.slides[slide_idx].shapes:
                if shape.shape_id == shape_id:
                    return shape

    def EditHeader(self):
        """Insert header description (title, contact info, photo, etc.).

        Shape name                  Slide number         Shape ID
        ----------                  -----------     -------------------
        Report title                     1          35 (PM2.5),  9 (O3)
        Report title                     2          21 (PM2.5), 21 (O3)
        Report title                     3          13 (PM2.5), 15 (O3)
        Deployment, contact info         1          34 (PM2.5), 33 (O3)
        Deployment, contact info         2          20 (PM2.5), 22 (O3)
        Deployment, contact info         3          12 (PM2.5), 16 (O3)
        Photo placeholder                1           2 (PM2.5),  2 (O3)
        Photo placeholder                2           4 (PM2.5),  3 (O3)
        Photo placeholder                3          14 (PM2.5),  3 (O3)

        Returns:
            None

        """
        # Title location for header
        title_left = 0.98
        title_top = 0.51
        text_vspace = 0
        # Get pptx text shape for modifying cells
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            title = self.GetShape(slide_idx, shape_loc=(title_left, title_top))

            # Set evaluation report title line 1: Parameter and test type
            title_line_1 = title.text_frame.paragraphs[0]
            run1 = title_line_1.add_run()
            run1.text = 'Testing Report - '
            param_baseline = title_line_1.add_run()
            param_baseline.text = self.param.format_baseline
            param_subscript = title_line_1.add_run()
            param_subscript.text = self.param.format_subscript
            font = param_subscript.font
            self.SetSubscript(font)
            run4 = title_line_1.add_run()
            run4.text = ' Base Testing'

            self.FormatText(title_line_1, alignment='left',
                            font_name='Calibri', font_size=30)
            title_line_1_font = title_line_1.font
            title_line_1_font.line_spacing = ppt.util.Pt(text_vspace)

            # Set evaluation report title line 2: Sensor name
            title_text_obj = title.text_frame.add_paragraph()
            title_text_obj.text = self.fmt_sensor_name
            self.FormatText(title_text_obj, alignment='left',
                            font_name='Calibri Light', font_size=30)
            title_font_obj = title_text_obj.font
            title_font_obj.line_spacing = ppt.util.Pt(text_vspace)

        # Header testing information location
        tester_left = 7.80
        tester_top = 0.36
        text_vspace = 10
        # Get pptx text shape for deployment and contact information
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            tester_info = self.GetShape(slide_idx, shape_loc=(tester_left,
                                                              tester_top))

            # Set deployment number
            tester_text = tester_info.text_frame.paragraphs[0]
            tester_text.text = self.testing_org['Deployment name']
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri', font_size=20, bold=True)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Set testing organization name (line 2)
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = self.testing_org['Org name'][0]
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Set testing organization name (line 2)
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = self.testing_org['Org name'][1]
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add contact email
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = (self.testing_org['Contact email'])
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add contact phone number
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = ('      ' + self.testing_org['Contact phone'])
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

            # Add current month to header
            month_year = dt.datetime.now().strftime('%B %Y')
            tester_text = tester_info.text_frame.add_paragraph()
            tester_text.text = month_year
            self.FormatText(tester_text, alignment='left',
                            font_name='Calibri Light', font_size=20)
            tester_text.line_spacing = ppt.util.Pt(text_vspace)

        # Edit header photo (edits the image placeholder)
        pic_left = 12.56
        pic_top = 0.29
        # Get pptx text shape for deployment and contact information
        for slide_n in np.arange(1, len(self.rpt.slides) + 1):
            slide_idx = int(slide_n) - 1
            pic = self.GetShape(slide_idx, shape_loc=(pic_left, pic_top))
            pic_path = self.figure_path + '\\deployment\\' + \
                self.name + '.png'
            if not os.path.exists(pic_path):
                print('No deployment picture found at', pic_path)
                placeholder_path = os.path.join(__file__,
                                                '../templates',
                                                'placeholder_image.png')
                try:
                    pic.insert_picture(placeholder_path)
                except AttributeError:
                    pass

            else:
                pic.insert_picture(pic_path)

    def EditSiteTable(self):
        """Add details to testing organzation and site info table (page 1).

        Table name                  Table ID
        ----------                  --------
        Testing org, site info         18


        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        shape = self.GetShape(slide_idx=0, shape_id=18)

        # ------------- Cell 1: Testing organization information --------------
        cell = shape.table.cell(1, 1)

        # Add organization name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = (self.testing_org['Org name'][0] + ' - ' +
                         self.testing_org['Org name'][1])
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=14)

        # Add contact link
        text_obj = cell.text_frame.add_paragraph()
        run = text_obj.add_run()
        run.text = self.testing_org['Website']['website name']
        link = self.testing_org['Website']['website link']
        hlink = run.hyperlink
        hlink.address = link
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=11)

        # Add contact email, phone number
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = (self.testing_org['Contact email'] + '\n' + '      ' +
                         self.testing_org['Contact phone'])
        self.FormatText(text_obj, alignment='left', font_name='Calibri',
                        font_size=11)

        # ----------- Cell 2: Testing location information ----------------
        cell = shape.table.cell(2, 1)

        # Add site name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.testing_loc['Site name']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # Add site address
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = self.testing_loc['Site address']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=11)

        # Add site lat long
        text_obj = cell.text_frame.add_paragraph()
        text_obj.text = (self.testing_loc['Site lat'] + ', ' +
                         self.testing_loc['Site long'])
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=11)

        # --------------- Cell 2: Testing location AQS ID ---------------------
        cell = shape.table.cell(3, 1)

        # Add site name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.testing_loc['Site AQS ID']
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # ------------------ Cell 3: Sampling timeframe -----------------------
        cell = shape.table.cell(4, 1)

        # Add sampling timeframe for each group in deployment
        self.eval_grps = list(self.tframe.keys())
        for i, grp in enumerate(self.tframe):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = self.eval_grps[i]
            grp_tframe = self.tframe[grp_name]

            if len(self.eval_grps) == 1:
                text_obj.text = grp_tframe
            else:
                text_obj.text = grp_name + ': ' + grp_tframe
            self.FormatText(text_obj, alignment='center', font_name='Calibri',
                            font_size=11)

    def EditSensorTable(self):
        """Add information to sensor information table (page 1).

        Table name                  Table ID
        ----------                  --------
        Sensor info             49 (PM2.5), 30 (O3)

        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        if self.n_avg_intervals == 2:
            shape = self.GetShape(slide_idx=0, shape_id=49)
        if self.n_avg_intervals == 1:
            shape = self.GetShape(slide_idx=0, shape_id=30)

        # Populate list with configured sensor recording interval(s)
        rec_interval = []
        for grp in self.deploy_dict['Deployment Groups']:
            grp_dict = self.deploy_dict['Deployment Groups'][grp]
            grp_sensors = grp_dict['sensors']
            for sensor in grp_sensors:
                interval = grp_sensors[sensor]['recording_interval']
                interval = interval.replace('.0', '')
                rec_interval.append(interval)
        rec_interval = list(set(rec_interval))
        rec_str = ', '.join(string for string in rec_interval)

        # ------------- Cell 1: Sensor manufacturer and model -----------------
        cell = shape.table.cell(1, 1)

        # Add sensor manufacturer, model name
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.fmt_sensor_name
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # --------------- Cell 3: Sensor recording interval -------------------
        cell = shape.table.cell(3, 1)

        # Add recording interval
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = rec_str
        self.FormatText(text_obj, alignment='center', font_name='Calibri',
                        font_size=14)

        # -------------------- Sensor Serial ID cells -------------------------
        rows = [4, 5, 6]  # row indicies for sensor serial cells
        cols = [1, 3, 4]  # column indicies for sensor serial cells
        for i, iloc in enumerate(rows):
            for j, jloc in enumerate(cols):
                cell_n = 3*i + (j + 1)

                # Only fill in ID info for number of sensors in evaluation.
                # If number of sensors is, say, 8, then the last cell for the
                # 3x3 grid for serial numbers is left empty.
                if len(self.serial_grp_dict) < cell_n:
                    pass
                else:
                    cell = shape.table.cell(iloc, jloc)

                    serial_idx = cell_n - 1
                    sensor_serial = list(self.serial_grp_dict.keys()
                                         )[serial_idx]
                    sensor_grp = list(self.serial_grp_dict.values()
                                      )[serial_idx]

                    # Add group number if multiple groups
                    if len(self.eval_grps) > 1:
                        grp_obj = cell.text_frame.paragraphs[0]
                        grp_obj.text = sensor_grp
                        self.FormatText(grp_obj, alignment='center',
                                        font_name='Calibri', font_size=10.5)

                        # Add sensor serial ID
                        serial_obj = cell.text_frame.add_paragraph()
                    else:
                        serial_obj = cell.text_frame.paragraphs[0]

                    serial_obj.text = sensor_serial
                    self.FormatText(serial_obj, alignment='center',
                                    font_name='Calibri', font_size=10.5)

        # Check for issues with deployment
        for i, grp in enumerate(self.deploy_dict['Deployment Groups']):
            deploy_grp_data = (self.deploy_dict['Deployment Groups']
                                               [grp]['sensors'])
            grp_sensors = list(deploy_grp_data.keys())
            grp_status = [deploy_grp_data[str(j)]['deploy_issues']
                          for j in grp_sensors]

            cell = shape.table.cell(7, 2)
            if list(set(grp_status)) == ['False']:
                if i == 0:
                    textobj = cell.text_frame.paragraphs[0]
                else:
                    textobj = cell.text_frame.add_paragraph()

                if len(self.eval_grps) == 1:
                    textobj.text = 'No Issues'
                else:
                    textobj.text = grp + ': No Issues'
                self.FormatText(textobj, alignment='center',
                                font_name='Calibri', font_size=12)
            else:
                if i == 0:
                    textobj = cell.text_frame.paragraphs[0]
                else:
                    textobj = cell.text_frame.add_paragraph()

                if len(self.eval_grps) == 1:
                    textobj.text = 'Issues with deployment'
                else:
                    textobj.text = grp + ': Issues with deployment'
                self.FormatText(textobj, alignment='center',
                                font_name='Calibri', font_size=12)

    def EditRefTable(self):
        """Add details to reference information table.

        Table name                  Table ID
        ----------                  --------
        Reference info                 51

        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        shape = self.GetShape(slide_idx=0, shape_id=51)

        # ------------- Cell 1: Sensor manufacturer and model -----------------
        cell = shape.table.cell(1, 1)

        # Add reference manufacturer and model
        text_obj = cell.text_frame.paragraphs[0]
        text_obj.text = self.ref_name
        self.FormatText(text_obj, alignment='center',
                        font_name='Calibri', font_size=14)

    def EditRefConcTable(self):
        """Add reference concentration tabular statistics (page 1).

        Located in different boxes based on the evaluation parameter type.

        Scatter plots box (PM2.5 only)
        Table name                  Table ID
        ----------                  --------
        Reference conc info          75 (PM25)

        Time series box (O3 only)
        Table name                  Table ID
        ----------                  --------
        Reference conc info          56 (O3)

        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        if self.n_avg_intervals == 2:
            shape = self.GetShape(slide_idx=0, shape_id=75)
        if self.n_avg_intervals == 1:
            shape = self.GetShape(slide_idx=0, shape_id=56)

        grp_info = self.deploy_dict['Deployment Groups']
        # reference concentration range
        self.refconc = {}
        for grp in list(grp_info.keys()):
            ref = 'Reference'
            try:
                ref_hmin = grp_info[grp][self._param_name][ref][
                                                            'conc_min_1-hour']
                ref_hmax = grp_info[grp][self._param_name][ref][
                                                            'conc_max_1-hour']
                ref_dmin = grp_info[grp][self._param_name][ref][
                                                            'conc_min_24-hour']
                ref_dmax = grp_info[grp][self._param_name][ref][
                                                            'conc_max_24-hour']

                self.refconc[grp] = \
                    '{0:3.1f}-{1:3.1f} (1-hr), '\
                    '{2:3.1f}-{3:3.1f} (24-hr)'.format(ref_hmin, ref_hmax,
                                                       ref_dmin, ref_dmax)

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # Number of periods reference exceeded concentration target
        if self.n_avg_intervals == 2:
            exceed_str = 'n_exceed_conc_goal_24-hour'
        if self.n_avg_intervals == 1:
            exceed_str = 'n_exceed_conc_goal_1-hour'

        self.refexceed = {}
        for grp in list(grp_info.keys()):
            try:
                self.refexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][self._param_name]['Reference'][exceed_str]
                       )

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # ------------- Cell 1: Range of ref concentrations -----------------
        cell = shape.table.cell(0, 1)
        for i, grp in enumerate(self.refconc):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.refconc.keys())[i]
            grp_refconc = self.refconc[grp_name]

            if len(self.eval_grps) == 1:
                text_obj.text = grp_refconc
            else:
                text_obj.text = grp_name + ': ' + grp_refconc
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=9)

        # ------------- Cell 2: N periods meeting conc. target-----------------
        cell = shape.table.cell(1, 1)
        for i, grp in enumerate(self.refexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.refexceed.keys())[i]
            grp_refexceed = self.refexceed[grp_name]
            if len(self.eval_grps) == 1:
                text_obj.text = grp_refexceed
            else:
                text_obj.text = grp_name + ': ' + grp_refexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=9)

    def EditMetCondTable(self):
        """Add meteorological conditions table (page 1).

        Table name                     Table ID
        ----------                     --------
        N outside target criteria   45 (O3), 74 (PM25)

        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        if self.n_avg_intervals == 2:
            shape = self.GetShape(slide_idx=0, shape_id=74)
        if self.n_avg_intervals == 1:
            shape = self.GetShape(slide_idx=0, shape_id=32)

        grp_info = self.deploy_dict['Deployment Groups']

        # Number of 24-hr periods temp exceeded target criteria
        self.tempexceed = {}
        exceed_str = 'n_exceed_target_criteria_24-hour'
        met_conds = 'Meteorological Conditions'  # abbreviating
        temp = 'Temperature'  # abbreviating
        for grp in list(grp_info.keys()):
            try:
                self.tempexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][met_conds][temp][exceed_str])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # Number of 24-hr periods temp exceeded target criteria
        self.rhexceed = {}
        exceed_str = 'n_exceed_target_criteria_24-hour'
        met_conds = 'Meteorological Conditions'  # abbreviating
        rh = 'Relative Humidity'  # abbreviating
        for grp in list(grp_info.keys()):
            try:
                self.rhexceed[grp] = \
                    '{0:d}'.format(
                       grp_info[grp][met_conds][rh][exceed_str])

            # Raise when attributes are 'none' likely due to no data
            except TypeError:
                pass

        # ----------- Cell 1: N periods outside temp target criteria ----------
        cell = shape.table.cell(0, 1)
        for i, grp in enumerate(self.tempexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.tempexceed.keys())[i]
            grp_tempexceed = self.tempexceed[grp_name]

            if len(self.eval_grps) == 1:
                text_obj.text = grp_tempexceed
            else:
                text_obj.text = grp_name + ': ' + grp_tempexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=10)

        # ----------- Cell 2: N periods outside rh target criteria ------------
        cell = shape.table.cell(1, 1)
        for i, grp in enumerate(self.rhexceed):
            if i == 0:
                text_obj = cell.text_frame.paragraphs[0]
            else:
                text_obj = cell.text_frame.add_paragraph()
            grp_name = list(self.rhexceed.keys())[i]
            grp_rhexceed = self.rhexceed[grp_name]

            if len(self.eval_grps) == 1:
                text_obj.text = grp_rhexceed
            else:
                text_obj.text = grp_name + ': ' + grp_rhexceed
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=10)

    def EditMetInfTable(self):
        """Add meteorological influence table (page 1).

        Table name                  Table ID
        ----------                  --------
        N paired met conc vals   48 (O3), 76 (PM25)

        Returns:
            None.

        """
        # Get pptx table shape for modifying cells
        if self.n_avg_intervals == 2:
            shape = self.GetShape(slide_idx=0, shape_id=76)
        if self.n_avg_intervals == 1:
            shape = self.GetShape(slide_idx=0, shape_id=33)

        grp_info = self.deploy_dict['Deployment Groups']

        # Number of 1-hr periods temp exceeded target criteria
        params = ['Temperature', 'Relative Humidity']
        met_conds = 'Meteorological Conditions'
        pair_str = 'n_measurement_pairs_1-hour'

        dic = {}

        for i, param in enumerate(params):
            dic[param] = {grp: grp_info[grp][met_conds][param][pair_str]
                          for grp in list(grp_info.keys())}

            for j, grp in enumerate(self.deploy_dict['Deployment Groups']):
                cell = shape.table.cell(i, 1)
                if j == 0:
                    text_obj = cell.text_frame.paragraphs[0]
                else:
                    text_obj = cell.text_frame.add_paragraph()
                grp_val = dic[param][grp]

                if len(self.eval_grps) == 1:
                    text_obj.text = str(round(grp_val))
                else:
                    text_obj.text = grp + ': ' + str(round(grp_val))
                self.FormatText(text_obj, alignment='center',
                                font_name='Calibri', font_size=10)

    def EditSensorRefTable(self, table):
        """Add sensor-reference tabular statistics (page 2).


        Args:
            table (TYPE): DESCRIPTION.

        Returns:
            None.

        """
        if self.n_avg_intervals == 2:
            span_dict = {'Bias and Linearity': [1, 6],
                         'Data Quality': [7, 10],
                         'R^2': [12, 13],
                         'Slope': [14, 15],
                         f'Intercept (b)\n({self.param.units})': [16, 17],
                         'Uptime (%)': [18, 19],
                         'N pairs': [20, 21]}
            table_categories = {'1': 'Bias and Linearity',
                                '7': 'Data Quality'}
            metrics = {'12': 'R^2',
                       '14': 'Slope',
                       '16': f'Intercept\n({self.param.units})',
                       '18': 'Uptime\n(%)',
                       '20': 'Number of paired\nsensor and '
                             'FRM/FEM\nconcentration pairs'}
            avg_intervals = {'23': '1-Hour',
                             '24': '24-Hour',
                             '25': '1-Hour',
                             '26': '24-Hour',
                             '27': '1-Hour',
                             '28': '24-Hour',
                             '29': '1-Hour',
                             '30': '24-Hour',
                             '31': '1-Hour',
                             '32': '24-Hour'}

            rsqr_lbound = self.param.PerformanceTargets.get_metric(
                                      metrics['12'])['bounds'][0]
            slope_goal = self.param.PerformanceTargets.get_metric(
                                      metrics['14'])['goal']
            slope_tol = self.param.PerformanceTargets.get_metric(
                                      metrics['14'])['bounds'][1] - slope_goal
            intcpt_lbound = self.param.PerformanceTargets.get_metric(
                                    metrics['16'].split('\n')[0])['bounds'][0]
            intcpt_ubound = self.param.PerformanceTargets.get_metric(
                                    metrics['16'].split('\n')[0])['bounds'][1]

            metric_targets = {'33': 'Metric Target Range',
                              '34': '≥ {:3.2f}'.format(rsqr_lbound),
                              '35': '≥ {:3.2f}'.format(rsqr_lbound),
                              '36': '{:2.1f} ± {:3.2f}'.format(slope_goal,
                                                               slope_tol),
                              '37': '{:2.1f} ± {:3.2f}'.format(slope_goal,
                                                               slope_tol),
                              '38': '{:1.0f} ≤ b ≤ {:1.0f}'.format(
                                                              intcpt_lbound,
                                                              intcpt_ubound),
                              '39': '{:1.0f} ≤ b ≤ {:1.0f}'.format(
                                                              intcpt_lbound,
                                                              intcpt_ubound),
                              '40': '75%*',
                              '41': '75%*',
                              '42': '-',
                              '43': '-'}

        if self.n_avg_intervals == 1:
            span_dict = {'Bias and Linearity': [1, 3],
                         'Data Quality': [4, 5]}
            table_categories = {'1': 'Bias and Linearity',
                                '4': 'Data Quality'}
            metrics = {'7': 'R^2',
                       '8': 'Slope',
                       '9': f'Intercept\n({self.param.units})',
                       '10': 'Uptime\n(%)',
                       '11': 'Number of paired\nsensor and '
                             'reference\nconcentration pairs'}
            avg_intervals = {'13': '1-Hour',
                             '14': '1-Hour',
                             '15': '1-Hour',
                             '16': '1-Hour',
                             '17': '1-Hour'}

            rsqr_lbound = self.param.PerformanceTargets.get_metric(
                                      metrics['7'])['bounds'][0]
            slope_goal = self.param.PerformanceTargets.get_metric(
                                      metrics['8'])['goal']
            slope_tol = self.param.PerformanceTargets.get_metric(
                                      metrics['8'])['bounds'][1] - slope_goal
            intcpt_lbound = self.param.PerformanceTargets.get_metric(
                                    metrics['9'].split('\n')[0])['bounds'][0]
            intcpt_ubound = self.param.PerformanceTargets.get_metric(
                                    metrics['9'].split('\n')[0])['bounds'][1]

            metric_targets = {'18': 'Metric Target Range',
                              '19': '≥ {:3.2f}'.format(rsqr_lbound),
                              '20': '{:2.1f} ± {:3.2f}'.format(slope_goal,
                                                               slope_tol),
                              '21': '{:1.0f} ≤ b ≤ {:1.0f}'.format(
                                                              intcpt_lbound,
                                                              intcpt_ubound),
                              '22': '75%*',
                              '23': '-'}

        self.grp_stats = self.stats_df.dropna()

        cells = self.SetSpanningCells(table, span_dict)

        if self.n_avg_intervals == 1:
            c1_row, c2_row, c3_row, c4_row, c5_row = 0, 0, 0, 0, 0
            h_stats = self.grp_stats.where(
                            self.grp_stats['Averaging Interval'] == '1-hour'
                            ).dropna().reset_index(drop=True)

        if self.n_avg_intervals == 2:
            (c1_row, c2_row, c3_row, c4_row, c5_row, c6_row, c7_row, c8_row,
             c9_row, c10_row) = (0, 0, 0, 0, 0, 0, 0, 0, 0, 0)
            h_stats = self.grp_stats.where(
                        self.grp_stats['Averaging Interval'] == '1-hour'
                        ).dropna().reset_index(drop=True)
            d_stats = self.grp_stats.where(
                        self.grp_stats['Averaging Interval'] == '24-hour'
                        ).dropna().reset_index(drop=True)

        # Sensor group dictionary, reset group keys for indexing
        grp_df = (self.deploy_dict['Deployment Groups']
                                  [self.grp_name]['sensors'])
        grp_df = {i: grp_df[k] for i, k in enumerate(sorted(grp_df.keys()))}

        # Lists for storing metric values to compute mean across sensors
        rsqr_h_vals, rsqr_d_vals = [], []
        slope_h_vals, slope_d_vals = [], []
        intcpt_h_vals, intcpt_d_vals = [], []
        self.uptime_h_vals, self.uptime_d_vals = [], []
        n_h_vals, n_d_vals = [], []

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]

            n_header_rows = 4
            if self.n_avg_intervals == 1:
                datacols = 6
            if self.n_avg_intervals == 2:
                datacols = 11
            datacell0 = datacols*n_header_rows
            datacellend = datacell0 + self.grp_n_sensors*datacols

            # Add table header info (rows 1-4)
            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                metric = metrics[str(i)]

                # Format the superscript for R^2
                if metric == 'R^2':
                    baseline = text_obj.add_run()
                    baseline.text = 'R'
                    superscript = text_obj.add_run()
                    superscript.text = '2'
                    font = superscript.font
                    self.SetSuperscript(font)
                # Format superscript for intercept if units in ug/m^3
                elif metric == 'Intercept\n(μg/m^3)':
                    lineone = text_obj.add_run()
                    lineone.text = 'Intercept\n'
                    linetwo_baseline_one = text_obj.add_run()
                    linetwo_baseline_one.text = '(μg/m'
                    superscript = text_obj.add_run()
                    superscript.text = '3'
                    font = superscript.font
                    self.SetSuperscript(font)
                    linetwo_baseline_two = text_obj.add_run()
                    linetwo_baseline_two.text = ')'
                else:
                    text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if i == datacell0 + self.grp_n_sensors*datacols:
                # Row of avg values at bottom of table
                text_obj.text = 'Mean'

            # Add Sensor Serial IDs to column 1
            if (i % datacols == 0 and i >= datacell0
               and c1_row < self.grp_n_sensors):
                try:
                    text_obj.text = 'Sensor ' + grp_df[c1_row]['serial_id']
                except KeyError:
                    pass

            # Decision tree for adding sensor data to table
            if (i > datacell0 and i < datacell0 + (datacols)*self.grp_n_sensors
               and i % datacols != 0):
                j = i - (datacell0 + 1)

                # Metric column 1
                if j % datacols == 0:
                    try:
                        # 1-hr R^2
                        val = h_stats.loc[c1_row, 'R$^2$']
                        text_obj.text = format(val, '3.2f')
                        rsqr_h_vals.append(val)
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c1_row += 1

                # Metric column 2
                if (j-1) % datacols == 0:
                    try:
                        if self.n_avg_intervals == 1:
                            # 1-hr slope
                            val = h_stats.loc[c2_row, 'Slope']
                            slope_h_vals.append(val)
                        if self.n_avg_intervals == 2:
                            # 24-hr R^2
                            val = d_stats.loc[c2_row, 'R$^2$']
                            rsqr_d_vals.append(val)

                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c2_row += 1

                # Metric column 3
                if (j-2) % datacols == 0:
                    try:
                        if self.n_avg_intervals == 1:
                            # 1-hr Intercept
                            val = h_stats.loc[c3_row, 'Intercept']
                            intcpt_h_vals.append(val)
                        if self.n_avg_intervals == 2:
                            # 1-hr Slope
                            val = h_stats.loc[c3_row, 'Slope']
                            slope_h_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c3_row += 1

                # Metric column 4
                if (j-3) % datacols == 0:
                    # insert Uptime
                    try:
                        if self.n_avg_intervals == 1:
                            # 1-hr uptime
                            val = grp_df[c4_row]['uptime_1-hour']
                            self.uptime_h_vals.append(val)
                        if self.n_avg_intervals == 2:
                            # 24-hr Slope
                            val = d_stats.loc[c4_row, 'Slope']
                            slope_d_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c4_row += 1

                # Metric column 5
                if (j-4) % datacols == 0:
                    # insert N paired sensor/reference values
                    try:
                        if self.n_avg_intervals == 1:
                            # 1-hr N
                            val = h_stats.loc[c5_row, 'N']
                            n_h_vals.append(val)
                        if self.n_avg_intervals == 2:
                            # 1-hr Intercept
                            val = h_stats.loc[c5_row, 'Intercept']
                            intcpt_h_vals.append(val)
                        text_obj.text = format(val, '3.2f')
                    except KeyError:  # Condition for empty grp stats
                        pass
                    c5_row += 1

                # Columns 6-10 only in PM2.5 sensor-reference table
                if self.n_avg_intervals == 2:

                    # Metric column 6
                    if (j-5) % datacols == 0:
                        try:
                            # 24-hr Intercept
                            val = d_stats.loc[c6_row, 'Intercept']
                            intcpt_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c6_row += 1

                    # Metric column 7
                    if (j-6) % datacols == 0:
                        try:
                            # 1-hr Uptime
                            val = grp_df[c7_row]['uptime_1-hour']
                            self.uptime_h_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c7_row += 1

                    # Metric column 8
                    if (j-7) % datacols == 0:
                        try:
                            # 24-hr Uptime
                            val = grp_df[c8_row]['uptime_24-hour']
                            self.uptime_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c8_row += 1

                    # Metric column 9
                    if (j-8) % datacols == 0:
                        try:
                            # 1-hr N
                            val = h_stats.loc[c9_row, 'N']
                            n_h_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c9_row += 1

                    # Metric column 10
                    if (j-9) % datacols == 0:
                        try:
                            # 24-hr N
                            val = d_stats.loc[c10_row, 'N']
                            n_d_vals.append(val)
                            text_obj.text = format(val, '3.2f')
                        except KeyError:  # Condition for empty grp stats
                            pass
                        c10_row += 1

            if i > datacellend:
                j = i - (datacellend + 1)

                # Metric Mean column 1
                if j % datacols == 0:
                    # Mean 1-hr R^2
                    if self.n_avg_intervals == 1:
                        val = np.mean(rsqr_h_vals)
                        txt = self.CheckTargets(rsqr_h_vals, metric='R^2')
                    # Mean 1-hr R^2
                    if self.n_avg_intervals == 2:
                        val = np.mean(rsqr_h_vals)
                        txt = self.CheckTargets(rsqr_h_vals, metric='R^2')

                # Metric Mean column 2
                if (j - 1) % datacols == 0:
                    # Mean 1-hr Slope
                    if self.n_avg_intervals == 1:
                        val = np.mean(slope_h_vals)
                        txt = self.CheckTargets(slope_h_vals, metric='Slope')
                    # Mean 24-hr R^2
                    if self.n_avg_intervals == 2:
                        val = np.mean(rsqr_d_vals)
                        txt = self.CheckTargets(rsqr_d_vals, metric='R^2')

                # Metric Mean column 3
                if (j - 2) % datacols == 0:
                    # Mean 1-hr Intercept
                    if self.n_avg_intervals == 1:
                        val = np.mean(intcpt_h_vals)
                        txt = self.CheckTargets(intcpt_h_vals,
                                                metric='Intercept')
                    # Mean 1-hr Slope
                    if self.n_avg_intervals == 2:
                        val = np.mean(slope_h_vals)
                        txt = self.CheckTargets(slope_h_vals, metric='Slope')

                # Metric Mean column 4
                if (j - 3) % datacols == 0:
                    # Mean 1-hr Uptime
                    if self.n_avg_intervals == 1:
                        val = np.mean(self.uptime_h_vals)
                        txt = self.CheckTargets(self.uptime_h_vals,
                                                metric='Uptime')
                    # Mean 24-hr Slope
                    if self.n_avg_intervals == 2:
                        val = np.mean(slope_d_vals)
                        txt = self.CheckTargets(slope_d_vals, metric='Slope')

                # Metric Mean column 5
                if (j - 4) % datacols == 0:
                    # Mean 1-hr N paired measurements
                    if self.n_avg_intervals == 1:
                        val = np.mean(n_h_vals)
                    # Mean 1-hr Intercept
                    if self.n_avg_intervals == 2:
                        val = np.mean(intcpt_h_vals)
                        txt = self.CheckTargets(intcpt_h_vals,
                                                metric='Intercept')

                if self.n_avg_intervals == 2:
                    # Metric Mean column 6
                    if (j - 5) % datacols == 0:
                        # Mean 24-hr Intercept
                        val = np.mean(intcpt_d_vals)
                        txt = self.CheckTargets(intcpt_d_vals,
                                                metric='Intercept')

                    # Metric Mean column 7
                    if (j - 6) % datacols == 0:
                        # Mean 1-hr Uptime
                        val = np.mean(self.uptime_h_vals)
                        txt = self.CheckTargets(self.uptime_h_vals,
                                                metric='Uptime')

                    # Metric Mean column 8
                    if (j - 7) % datacols == 0:
                        # Mean 24-hr Uptime
                        val = np.mean(self.uptime_d_vals)
                        txt = self.CheckTargets(self.uptime_d_vals,
                                                metric='Uptime')

                    # Metric Mean column 9
                    if (j - 8) % datacols == 0:
                        # Mean 1-hr N paired measurements
                        val = np.mean(n_h_vals)

                    # Metric Mean column 10
                    if (j - 9) % datacols == 0:
                        # Mean 24-hr N paired measurements
                        val = np.mean(n_d_vals)

                # Indicate number of sensors meeting performance metric target
                if txt is not None:
                    trgt_cell = cells[i - datacols*(self.grp_n_sensors + 2)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

                text_obj.text = format(val, '3.2f')

            # Configure text formatting
            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)
            txt = None

    def EditErrorTable(self, table):
        """Add error tabular statistics (page 2).


        Args:
            table (TYPE): DESCRIPTION.

        Returns:
            None.

        """

        error_stats = (self.deploy_dict['Deployment Groups']
                                       [self.grp_name]
                                       [self._param_name]
                                       ['Error'])

        for key, value in error_stats.items():
            if value is None:
                error_stats[key] = -999

        if self.n_avg_intervals == 2:
            datacols = 4
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Error': [1, 4],
                         f'RMSE\n({self.param.units})': [6, 7],
                         'NRMSE\n(%)': [8, 9]}
            table_categories = {'1': 'Error'}
            metrics = {'6': f'RMSE\n({self.param.units})',
                       '8': 'NRMSE\n(%)'}
            avg_intervals = {'11': '1-Hour',
                             '12': '24-Hour',
                             '13': '1-Hour',
                             '14': '24-Hour'}

            rmse_ubound = self.param.PerformanceTargets.get_metric(
                                      metrics['6'].split('\n')[0])['bounds'][1]
            nrmse_ubound = self.param.PerformanceTargets.get_metric(
                                      metrics['8'].split('\n')[0])['bounds'][1]

            metric_targets = {'15': 'Metric Target Range',
                              '16': '≤  {:2.1f}'.format(rmse_ubound),
                              '17': '≤ {:2.1f}'.format(rmse_ubound),
                              '18': '≤ {:3.1f}'.format(nrmse_ubound),
                              '19': '≤ {:3.1f}'.format(nrmse_ubound),
                              '20': 'Deployment Value'}
            metric_vals = {'21': format(error_stats['rmse_1-hour'], '3.2f'),
                           '22': format(error_stats['rmse_24-hour'], '3.2f'),
                           '23': format(error_stats['nrmse_1-hour'], '3.2f'),
                           '24': format(error_stats['nrmse_24-hour'], '3.2f')}
            for key, value in metric_vals.items():
                if value == '-999.00':
                    metric_vals[key] = ''

        if self.n_avg_intervals == 1:
            datacols = 1
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Error': [1, 1]}
            table_categories = {'1': 'Error'}
            metrics = {'3': f'RMSE\n({self.param.units})'}
            avg_intervals = {'5': '1-Hour'}

            rmse_ubound = self.param.PerformanceTargets.get_metric(
                                      metrics['3'].split('\n')[0])['bounds'][1]

            metric_targets = {'6': 'Metric Target Range',
                              '7': '≤ {:2.1f}'.format(rmse_ubound),
                              '8': 'Deployment Value'}
            metric_vals = {'9': format(error_stats['rmse_1-hour'], '3.2f')}
            for key, value in metric_vals.items():
                if value == '-999.00':
                    metric_vals[key] = ''

        cells = self.SetSpanningCells(table, span_dict)

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]

            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                metric = metrics[str(i)]

                # TODO: Possibly remove since using unicode chars?
                if metric == 'RMSE\n(μg/m^3)':
                    lineone = text_obj.add_run()
                    lineone.text = 'RMSE\n'
                    linetwo_baseline_one = text_obj.add_run()
                    linetwo_baseline_one.text = '(μg/m'
                    superscript = text_obj.add_run()
                    superscript.text = '3'
                    font = superscript.font
                    self.SetSuperscript(font)
                    linetwo_baseline_two = text_obj.add_run()
                    linetwo_baseline_two.text = ')'
                else:
                    text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if str(i) in metric_vals:
                text_obj.text = metric_vals[str(i)]

            # Metric column 1
            if i > headercellend:
                if (i - 1) % datacols == 0:
                    # Mean 1-hr RMSE
                    if self.n_avg_intervals == 1:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')
                    # Mean 1-hr RMSE
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')
                # Metric column 2
                if (i - 1) % datacols == 1:
                    # 1-hr NRMSE
                    if self.n_avg_intervals == 1:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='NRMSE')
                    # 24-hr RMSE
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='RMSE')
                # Metric column 3
                if self.n_avg_intervals == 2:
                    if (i - 1) % datacols == 2:
                        # 1-hr NRMSE
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='NRMSE')
                    # Metric column 4
                    if (i - 1) % datacols == 3:
                        # 24-hr NRMSE
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='NRMSE')

                # Indicate whether sensors meet performance metric target
                if txt is not None:
                    trgt_cell = cells[i - 2*(datacols + 1)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)

            txt = None

    def EditSensorSensorTable(self, table):
        """Add intersensor (sensor-sensor) precision tabular stats (page 2).


        Args:
            table (TYPE): DESCRIPTION.

        Returns:
            None.

        """

        grp_stats = (self.deploy_dict['Deployment Groups']
                                     [self.grp_name]
                                     [self._param_name]
                                     ['Precision'])

        if self.n_avg_intervals == 2:
            datacols = 8
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Precision (between collocated sensors)': [1, 4],
                         'Data Quality': [5, 8],
                         'CV\n(%)': [10, 11],
                         f'SD\n({self.param.units})': [12, 13],
                         'Uptime\n(%)': [14, 15],
                         'Number concurrent sensor values': [16, 17]}
            table_categories = {'1': 'Precision (between collocated sensors)',
                                '5': 'Data Quality'}
            metrics = {'10': 'CV\n(%)',
                       '12': f'SD\n({self.param.units})',
                       '14': 'Uptime\n(%)',
                       '16': 'Number of paired\nsensor and '
                             'reference\nconcentration pairs'}
            avg_intervals = {'19': '1-Hour',
                             '20': '24-Hour',
                             '21': '1-Hour',
                             '22': '24-Hour',
                             '23': '1-Hour',
                             '24': '24-Hour',
                             '25': '1-Hour',
                             '26': '24-Hour'}

            cv_ubound = self.param.PerformanceTargets.get_metric(
                            metrics['10'].split('\n')[0])['bounds'][1]
            sd_ubound = self.param.PerformanceTargets.get_metric(
                            metrics['12'].split('\n')[0])['bounds'][1]

            metric_targets = {'27': 'Metric Target Range',
                              '28': '≤ {:3.1f}'.format(cv_ubound),
                              '29': '≤ {:3.1f}'.format(cv_ubound),
                              '30': '≤ {:2.1f}'.format(sd_ubound),
                              '31': '≤ {:2.1f}'.format(sd_ubound),
                              '32': '75%*',
                              '33': '75%*',
                              '34': '-',
                              '35': '-',
                              '36': 'Deployment Value'}
            metric_vals = {'37': format(grp_stats['cv_1-hour'], '3.2f'),
                           '38': format(grp_stats['cv_24-hour'], '3.2f'),
                           '39': format(grp_stats['std_1-hour'], '3.2f'),
                           '40': format(grp_stats['std_24-hour'], '3.2f'),
                           '41': format(np.mean(self.uptime_h_vals), '3.2f'),
                           '42': format(np.mean(self.uptime_d_vals), '3.2f'),
                           '43': format(grp_stats['n_1-hour']),
                           '44': format(grp_stats['n_24-hour'])}

        if self.n_avg_intervals == 1:
            datacols = 4
            nheaderrows = 4
            headercellend = nheaderrows*(datacols + 1)
            span_dict = {'Precision (between collocated sensors)': [1, 2],
                         'Data Quality': [3, 4]}

            table_categories = {'1': 'Precision (between collocated sensors)',
                                '3': 'Data Quality'}
            metrics = {'6': 'CV\n(%)',
                       '7': f'SD\n({self.param.units})',
                       '8': 'Uptime\n(%)',
                       '9': 'Number of paired\nsensor and '
                            'reference\nconcentration pairs'}
            avg_intervals = {'11': '1-Hour',
                             '12': '1-Hour',
                             '13': '1-Hour',
                             '14': '1-Hour'}

            cv_ubound = self.param.PerformanceTargets.get_metric(
                            metrics['6'].split('\n')[0])['bounds'][1]
            sd_ubound = self.param.PerformanceTargets.get_metric(
                            metrics['7'].split('\n')[0])['bounds'][1]

            metric_targets = {'15': 'Metric Target Range',
                              '16': '≤ {:3.1f}'.format(cv_ubound),
                              '17': '≤ {:2.1f}'.format(sd_ubound),
                              '18': '75%*',
                              '19': '-',
                              '20': 'Deployment Value'}
            metric_vals = {'21': format(grp_stats['cv_1-hour']),
                           '22': format(grp_stats['std_1-hour']),
                           '23': format(np.mean(self.uptime_h_vals), '3.2f'),
                           '24': format(grp_stats['n_1-hour'])}

        cells = self.SetSpanningCells(table, span_dict)

        for i, cell in enumerate(cells):
            text_obj = cell.text_frame.paragraphs[0]
            if str(i) in table_categories:
                text_obj.text = table_categories[str(i)]
                font_obj = text_obj.font
                font_obj.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)
            if str(i) in metrics:
                metric = metrics[str(i)]

                # TODO: Possibly remove since using unicode chars?
                if metric == 'SD\n(μg/m^3)':
                    lineone = text_obj.add_run()
                    lineone.text = 'SD\n'
                    linetwo_baseline_one = text_obj.add_run()
                    linetwo_baseline_one.text = '(μg/m'
                    superscript = text_obj.add_run()
                    superscript.text = '3'
                    font = superscript.font
                    self.SetSuperscript(font)
                    linetwo_baseline_two = text_obj.add_run()
                    linetwo_baseline_two.text = ')'
                else:
                    text_obj.text = metrics[str(i)]
            if str(i) in metric_targets:
                text_obj.text = metric_targets[str(i)]
            if str(i) in avg_intervals:
                text_obj.text = avg_intervals[str(i)]
            if str(i) in metric_vals:
                text_obj.text = metric_vals[str(i)]

            # Metric column 1
            if i > headercellend:
                j = i - headercellend
                if (j - 1) % datacols == 0:
                    # Mean 1-hr CV
                    if self.n_avg_intervals == 1:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                    # Mean 1-hr CV
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 2
                if (j - 1) % datacols == 1:
                    # 1-hr Std dev
                    if self.n_avg_intervals == 1:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='SD')
                    # 24-hr CV
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 3
                if (j - 1) % datacols == 2:
                    # 1-hr Uptime
                    if self.n_avg_intervals == 1:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')
                    # 1-hr Std dev
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='CV')
                # Metric column 4
                if (j - 1) % datacols == 3:
                    # 24-hr Std dev
                    if self.n_avg_intervals == 2:
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='SD')

                if self.n_avg_intervals == 2:
                    # Metric column 5
                    if (j - 1) % datacols == 4:
                        # 1-hr uptime
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')
                    # Metric column 6
                    if (j - 1) % datacols == 5:
                        # 24-hr uptime
                        val = float(metric_vals[str(i)])
                        txt = self.CheckTargets(val, metric='Uptime')

                # Indicate whether sensors meet performance metric target
                if txt is not None:
                    trgt_cell = cells[i - 2*(datacols + 1)]
                    trgt_cell_text = trgt_cell.text_frame.add_paragraph()
                    trgt_cell_text.text = txt
                    self.FormatText(trgt_cell_text, alignment='center',
                                    font_name='Calibri Light', font_size=12)

            self.FormatText(text_obj, alignment='center',
                            font_name='Calibri', font_size=14)

            txt = None

    def SetSpanningCells(self, table, span_dict):
        """Merge tabular cells to form cells spanning multiple rows/columns.


        Args:
            table (TYPE): DESCRIPTION.
            span_dict (TYPE): DESCRIPTION.

        Returns:
            cells (TYPE): DESCRIPTION.

        """

        cells = [cell for cell in table.iter_cells()]

        if span_dict is not None:
            for entry in span_dict:
                merge_start_idx = span_dict[entry][0]
                merge_end_idx = span_dict[entry][1]

                merge_start_cell = cells[merge_start_idx]
                merge_end_cell = cells[merge_end_idx]
                merge_start_cell.merge(merge_end_cell)

        return cells

    def EditTabularStats(self):
        """Wrapper for constructing tables and adding entries (page 2).

        Returns:
            None.

        """
        self.n_grps = len(set(self.serial_grp_dict.values()))
        self.n_sensors = len(set(self.serial_grp_dict))

        # Use slide layout for generating additional slides
        tabular_layout_idx = 1
        tabular_layout = self.rpt.slide_layouts[tabular_layout_idx]

        legend_txt = {
                'line1': 'Device-specific metrics (computed for each '
                         'sensor in evaluation)',
                'line2': '○○○	Metric value for none of devices tested '
                         'falls within the target range',
                'line3': '●○○	Metric value for one of devices tested '
                         'falls within the target range',
                'line4': '●●○   Metric value for two of devices tested '
                         'falls within the target range',
                'line5': '●●●   Metric value for three of devices tested '
                         'falls within the target range',
                'line6': '',
                'line7': 'Single-valued metrics '
                         '(computed via entire evaluation dataset)',
                'line8': '○  Indicates that the metric value is not '
                         'within the target range',
                'line9': '●  Indicates that the metric value is within '
                         'the target range',

                }

        # List of unique testing groups
        grps = sorted(list(set(self.serial_grp_dict.values())))

        for grp_n, grp_name in enumerate(grps, 1):
            self.grp_n_sensors = list(
                        self.serial_grp_dict.values()).count(grp_name)
            self.grp_name = grp_name

            # Create new tabular stats page
            tabular_slide = self.rpt.slides.add_slide(
                                                tabular_layout)

            # Sensor Reference Table
            sr_frame, sr_table = self.ConstructTable(
                                        tabular_slide,
                                        table_type='sensor_reference')
            self.EditSensorRefTable(sr_table)

            # Error Table
            e_frame, e_table = self.ConstructTable(
                                        tabular_slide,
                                        table_type='error')
            self.EditErrorTable(e_table)

            # Intersensor (sensor-sensor) table
            ss_frame, ss_table = self.ConstructTable(
                                        tabular_slide,
                                        table_type='sensor_sensor')
            self.EditSensorSensorTable(ss_table)

            # Adjust sensor-reference table vertical position
            EMU_to_in = 1/914400.0
            sr_top = 4.03  # in
            sr_height = EMU_to_in*sr_frame.height
            sr_frame.top = ppt.util.Inches(sr_top)

            # Adjust error table vertical position
            e_t = 4.03 + sr_height + 0.4  # in
            e_h = EMU_to_in*e_frame.height
            e_frame.top = ppt.util.Inches(e_t)

            # Adjust Sensor-Sensor table vertical position
            ss_top = e_t + e_h + 2*0.4  # in
            ss_frame.top = ppt.util.Inches(ss_top)

            # Add section header
            tabular_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(2.81),  # top
                                ppt.util.Inches(3.16),  # width
                                ppt.util.Inches(0.47))  # height
            tabular_header_obj = tabular_header.text_frame.paragraphs[0]

            tabular_header_obj.text = 'Tabular Statistics'
            if self.n_grps > 1:
                tabular_header_obj.text += ' - ' + self.grp_name

            self.FormatText(tabular_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=22)

            # Add Sensor-Reference section header text label
            sr_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(3.30),  # top
                                ppt.util.Inches(12.81),  # width
                                ppt.util.Inches(0.44))  # height
            sr_header_obj = sr_header.text_frame.paragraphs[0]
            sr_header_obj.text = 'Sensor-FRM/FEM Correlation'
            self.FormatText(sr_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=20)

            # Add Sensor-Sensor section header text label
            ss_header = tabular_slide.shapes.add_textbox(
                                ppt.util.Inches(0.82),  # left
                                ppt.util.Inches(e_t + e_h + 0.12),  # top
                                ppt.util.Inches(12.81),  # width
                                ppt.util.Inches(0.44))  # height
            ss_header_obj = ss_header.text_frame.paragraphs[0]
            ss_header_obj.text = 'Sensor-Sensor Precision'
            self.FormatText(ss_header_obj, alignment='left',
                            font_name='Calibri Light', font_size=20)

            connector1 = tabular_slide.shapes.add_connector(
                                ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                                ppt.util.Inches(0.85),  # start_x
                                ppt.util.Inches(3.30),  # start_y
                                ppt.util.Inches(16.24),  # end_x
                                ppt.util.Inches(3.30))  # end_y

            connector1.line.fill.solid()
            connector1.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171,
                                                                         171,
                                                                         171)

            connector2 = tabular_slide.shapes.add_connector(
                            ppt.enum.shapes.MSO_CONNECTOR.STRAIGHT,
                            ppt.util.Inches(0.85),  # start_x
                            ppt.util.Inches(e_t + e_h + 0.56),  # start_y
                            ppt.util.Inches(16.24),  # end_x
                            ppt.util.Inches(e_t + e_h + 0.56))  # end_y

            connector2.line.fill.solid()
            connector2.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(171,
                                                                         171,
                                                                         171)

            # Add Sensor-Sensor text label
            legend_h = 2.7
            legend = tabular_slide.shapes.add_shape(
                                ppt.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE,
                                ppt.util.Inches(8.83),  # left
                                ppt.util.Inches(e_t + 0.5*(e_h + - legend_h)),
                                ppt.util.Inches(6.69),  # width
                                ppt.util.Inches(legend_h))

            legend.fill.solid()
            legend.fill.fore_color.rgb = ppt.dml.color.RGBColor(236,
                                                                236,
                                                                240)

            legend.line.fill.solid()
            legend.line.fill.fore_color.rgb = ppt.dml.color.RGBColor(214,
                                                                     216,
                                                                     226)
            legend.line.width = ppt.util.Pt(2.5)

            legend_obj = legend.text_frame.paragraphs[0]
            legend_obj.text = legend_txt['line1']
            self.FormatText(legend_obj, alignment='left',
                            font_name='Calibri', font_size=16, bold=False)

            for i, line in enumerate(legend_txt):
                if i > 0:
                    legend_obj = legend.text_frame.add_paragraph()
                    legend_obj.text = legend_txt[line]

                font = legend_obj.font
                font.name = 'Calibri'
                font.Bold = False
                font.color.rgb = ppt.dml.color.RGBColor(0, 0, 0)

                if line == 'line6':
                    font.size = ppt.util.Pt(8)
                else:
                    font.size = ppt.util.Pt(15)

    def ConstructTable(self, slide, table_type='sensor_reference'):
        """Select and construct tables on report page 2.

        Presets are set for constructing each table type (number of rows and
        columns, dimensions of tables, shading of cells and fill color, etc.)

        Args:
            slide (pptx slide object):
                The report slide on which the tabular statistics will be
                placed. This will likely be slide #2 (i.e., self.rpt.slides[1])
            table_type (str): {'sensor_reference', 'error', 'sensor_sensor'}
                Name of the type of table to construct.

        Returns:
            frame:
                pptx GraphicFrame object that the table is contained within
            table:
                pptx table shape formatted for the selected table type

        """

        cell_margin = 0.001
        table_spacing = 0.4
        shapes = slide.shapes

        # Create Sensor-Reference Correlation Table ---------------------------
        if table_type == 'sensor_reference':
            if self.n_avg_intervals == 2:
                nrows = 5 + self.grp_n_sensors
                ncols = 11
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(14.4)
                height = ppt.util.Inches(3.07 + 0.45*self.grp_n_sensors)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69)
                # grey out cells at index (from l-r where top left is 0)
                greyed_cells = [42, 43]

            if self.n_avg_intervals == 1:
                nrows = 5 + self.grp_n_sensors
                ncols = 6
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(14.4)
                height = ppt.util.Inches(3.07 + 0.45*self.grp_n_sensors)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69)
                greyed_cells = [23]

        if table_type == 'error':
            if self.n_avg_intervals == 2:
                nrows = 5
                ncols = 5
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(7.2)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + table_spacing)
                greyed_cells = None

            if self.n_avg_intervals == 1:
                nrows = 5
                ncols = 2
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(7.2)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + table_spacing)
                greyed_cells = None

        if table_type == 'sensor_sensor':
            if self.n_avg_intervals == 2:
                nrows = 5
                ncols = 9
                col_width = ppt.util.Inches(1.2)
                width = ppt.util.Inches(12.0)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + 2*table_spacing + 3.23)
                greyed_cells = [34, 35]

            if self.n_avg_intervals == 1:
                nrows = 5
                ncols = 5
                col_width = ppt.util.Inches(2.4)
                width = ppt.util.Inches(12.0)
                height = ppt.util.Inches(3.23)
                left = ppt.util.Inches(1.32)
                top = ppt.util.Inches(3.69 + 7.83 + 2*table_spacing + 3.23)
                greyed_cells = [19]

        # Construct the table based on selected presets
        frame = shapes.add_table(nrows, ncols, left, top,
                                 width, height)
        table = frame.table
        table.horz_banding = False

        # Loop over the cells in the table, configure fill color, cell margins
        for cell_idx, cell in enumerate(table.iter_cells()):
            # Set cell border width
            self.SetCellBorder(cell)

            # Set dark blue cell fill color
            if cell_idx < 3*ncols + 1 or cell_idx % ncols == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ppt.dml.color.RGBColor(191,
                                                                  208,
                                                                  235)
            # Set transparent (background) fill for top 3 cells
            if cell_idx % ncols == 0 and cell_idx <= 2*ncols:
                cell.fill.background()

            # Grey out cells where no target range specified
            if greyed_cells is not None:
                if cell_idx in greyed_cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ppt.dml.color.RGBColor(214,
                                                                      216,
                                                                      226)
            # Set cell margins to near zero
            cell.margin_left = ppt.util.Inches(cell_margin)
            cell.margin_right = ppt.util.Inches(cell_margin)
            cell.margin_top = ppt.util.Inches(cell_margin)
            cell.margin_bottom = ppt.util.Inches(cell_margin)
            # Vertical position text in middle of cells
            cell.vertical_anchor = ppt.enum.text.MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = ppt.util.Pt(14)

        # Modify table column width
        table.columns[0].width = ppt.util.Inches(2.2)
        for col_idx in np.arange(1, len(table.columns)):
            table.columns[col_idx].width = col_width

        # Modify table header row height
        table.rows[0].height = ppt.util.Inches(0.47)  # Category
        table.rows[1].height = ppt.util.Inches(0.97)  # Metric label
        table.rows[2].height = ppt.util.Inches(0.72)  # Averaging interval
        for row_idx in np.arange(3, len(table.rows)):
            table.rows[row_idx].height = ppt.util.Inches(0.45)

        return frame, table

    def PrintpptxShapes(self, slide_number=1, shape_type='all'):
        """Diagnostic tool for indicating shape ids and locations on reporting
        template slides.

        Args:
            slide number (int):
                The number of the slide (starting at 1) for which shape ids and
                locations will be printed.
            shape_type (str) {'all', ...?}: # TODO: get type names
                The types of shapes on the slide to print out. 'all' will
                return all shapes regardless of type, however, selecting a
                particular type (e.g., 'table') will only return shapes on the
                page corresponding to the specified type.

        Returns:
            None

        """
        print("{:^6s}{:^18s}{:^12s}{:^10s}".format('ID', 'Type',
                                                   'Left loc', 'Top loc'))
        print("{:^6s}{:^18s}{:^12s}{:^10s}".format('', '', '(in)', '(in)'))
        print('{:^45s}'.format(45*'-'))
        if shape_type == 'all':
            for shape in self.rpt.slides[slide_number - 1].shapes:
                print("%4s %16s %10.2f %10.2f" %
                      (shape.shape_id, shape.shape_type,
                       shape.left.inches, shape.top.inches))
        else:
            for shape in self.rpt.slides[slide_number - 1].shapes:
                if str(shape.shape_type).startswith(shape_type):
                    print("%4s %16s %10.2f %10.2f" %
                          (shape.shape_id, shape.shape_type,
                           shape.left.inches, shape.top.inches))

    def SubElement(self, parent, tagname, **kwargs):
        """Modify xml entry to add a new attribute (element).

        Reference:
            Based on Steve Canny's code at the following link:
            https://groups.google.com/g/python-pptx/c/UTkdemIZICw

        """
        element = ppt.oxml.xmlchemy.OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def SetCellBorder(self, cell, border_color="ffffff", border_width='20000'):
        """Edit tabular cell boarder width, color.

        Reference:
            Based on Steve Canny's code at the following links:
            https://groups.google.com/g/python-pptx/c/UTkdemIZICw

            https://stackoverflow.com/questions/42610829/
            python-pptx-changing-table-style-or-adding-borders-to-cells

        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = self.SubElement(tcPr, lines, w=border_width,
                                 cap='flat', mpd='sng', algn='ctr')
            solidFill = self.SubElement(ln, 'a:solidFill')
            srgbClr = self.SubElement(solidFill, 'a:srgbClr',
                                      val=border_color)
            prstDash = self.SubElement(ln, 'a:prstDash',
                                       val='solid')
            round_ = self.SubElement(ln, 'a:round')
            headEnd = self.SubElement(ln, 'a:headEnd',
                                      type='none', w='med', len='med')
            tailEnd = self.SubElement(ln, 'a:tailEnd',
                                      type='none', w='med', len='med')

    def SetSubscript(self, font):
        """Workaround for making font object text subscript (not included in
        python-pptx as of v0.6.19)

        Reference:
            Code via Martin Packer:
            https://stackoverflow.com/questions/61329224/how-do-i-add-
            superscript-subscript-text-to-powerpoint-using-python-pptx

        """
        font._element.set('baseline', '-25000')

    def SetSuperscript(self, font):
        """Workaround for making font object text superscript (not included in
        python-pptx as of v0.6.19)

        Reference:
            Code via Martin Packer:
            https://stackoverflow.com/questions/61329224/how-do-i-add-
            superscript-subscript-text-to-powerpoint-using-python-pptx

        """
        font._element.set('baseline', '30000')

    def MoveSlide(self, slides, slide, new_idx):
        """Move the supplemnental info table to the last slide position.

        Reference:
            Code via github user Amazinzay (Feb 17 2021):
            https://github.com/scanny/python-pptx/issues/68

        """
        slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])

    def AddSlideNumbers(self):
        """Add slide numbers to slides generated during report construction.

        For some reason, the python pptx module can't assign the footer page
        number to slides that are created by the library. While slides that
        are imported via the template (the first and last page of the report)
        have page number placeholders already assigned, the pptx library doesnt
        do this without explicity copying and pasting the page number
        placeholder from the layout to the slides that are created by the
        module.

        Reference:
            This code follows the basic outline Steve Canny (scanny) suggests
            in response to this GitHub post:
            https://github.com/scanny/python-pptx/issues/64

        """
        layout = self.rpt.slide_layouts[1]
        placeholders = layout.placeholders
        for i, placeholder in enumerate(placeholders):
            if layout.placeholders[i].name == 'Slide Number Placeholder':
                break

        # add slide numbers to the 2nd through the 2nd to last slide
        # (slides that are generated by PerformanceReport)
        for idx in np.arange(1, len(self.rpt.slides)-1, 1):
            slide = self.rpt.slides[idx]
            slide.shapes.clone_placeholder(placeholder)
            slide_placeholder = slide.shapes[-1]
            slide_placeholder.text = str(idx + 1)

    def FormatText(self, text_obj, alignment='center', font_name='Calibri',
                   font_size=24, bold=False, italic=False):
        """Set text attributes (font, size, bold, italic, alignment).

        Args:
            text_obj (TYPE): DESCRIPTION.
            alignment (TYPE, optional): DESCRIPTION. Defaults to 'center'.
            font_name (TYPE, optional): DESCRIPTION. Defaults to 'Calibri'.
            font_size (TYPE, optional): DESCRIPTION. Defaults to 24.
            bold (TYPE, optional): DESCRIPTION. Defaults to False.
            italic (TYPE, optional): DESCRIPTION. Defaults to False.

        Returns:
            None.

        """

        if alignment == 'center':
            text_obj.alignment = ppt.enum.text.PP_ALIGN.CENTER
        if alignment == 'left':
            text_obj.alignment = ppt.enum.text.PP_ALIGN.LEFT

        font_obj = text_obj.font
        font_obj.name = font_name
        font_obj.size = ppt.util.Pt(font_size)
        font_obj.bold = bold
        font_obj.italic = italic

    def CheckTargets(self, metric_vals, metric):
        """Evaluate how many sensors met a metric target, return textual
        depiction.

        For a passed metric name 'metric', determine the number of sensors
        with metric values within the specified metric target range.

        Example:
            Say the 'metric' argument is 'CV' and the 'metric_vals' argument is
            [20.2, 43.6, 26.5] (values are percentages). Given that the target
            range for 'CV' is from 0% to 30%, two our of three sensors fall
            within the target range. Textually, this can be represented by
            a series of three dots, where two dots are closed and one is empty.

            Text returned by CheckTargets():
                '●●○'

        Args:
            metric_vals (TYPE): DESCRIPTION.
            metric (TYPE): DESCRIPTION.

        Returns:
            text (TYPE): DESCRIPTION.

        """
        # Place float / int values (single-valued intersensor stats) into
        # list for parsing
        if type(metric_vals) != list:
            metric_vals = [metric_vals]

        if metric != 'Uptime':
            metric_info = self.param.PerformanceTargets.get_metric(metric)
            metric_bounds = metric_info['bounds']
            metric_min, metric_max = metric_bounds
        else:
            Uptime_min, Uptime_max = 75, 100
            metric_min, metric_max = Uptime_min, Uptime_max

        # Number of sensors meeting target
        n_sensors = len(metric_vals)
        n_passed = sum(metric_min <= val <= metric_max for val
                       in metric_vals)

        open_dot = '○'
        closed_dot = '●'

        if n_sensors != 0:
            pcnt_passed = 100*(n_passed / n_sensors)

        if metric in ['R^2', 'Slope', 'Intercept', 'Uptime']:
            text = n_passed*closed_dot + (n_sensors - n_passed)*open_dot
            if n_sensors == 1 and metric == 'Uptime':
                if pcnt_passed == 100:
                    text = closed_dot
                else:
                    text = open_dot

        if metric in ['CV', 'SD', 'RMSE', 'NRMSE']:
            if pcnt_passed == 100:
                text = closed_dot
            elif n_sensors > 0:
                text = open_dot
            else:
                # Metric values empty?
                text = ''

        return text

    def CreateReport(self):
        """Wrapper for running the various methods that construct reports.

        Existing figures are assumed to have been created on the same day of
        class instantiation. If a figure filename is not found, sensor data
        are loaded via the SensorEvaluation class and the figure is generated.

        Returns:
            None.

        """
        print('Creating Testing Report for', self.name)

        ref_source = self.hourly_ref_df.Data_Source.mode()[0]

        MET_DATA = True
        if self.met_hourly_ref_df.select_dtypes(exclude='object').dropna(how='all', axis=1).empty:
            print('')
            warnings.warn('Warning: No Meteorological data for reference'
                          ' source. Met plots will not be added to the report')
            MET_DATA = False

        # Set figure positions
        self.FigPositions()

        print('..Adding figures to report')
        # Add figures to report
        self.AddSingleScatterPlot()
        self.AddTimeseriesPlot()
        self.AddMetricsPlot()

        print('..Adding tabular data')
        # Modify report tables
        self.EditSiteTable()
        self.EditSensorTable()
        self.EditRefTable()
        self.EditRefConcTable()
        self.EditTabularStats()

        # Reference dependent details. Only add met plots if reference data
        # come from AirNowTech or another source where met data are provided
        if MET_DATA:
            self.AddMetDistPlot()
            self.AddMetInflPlot()
            self.EditMetCondTable()
            self.EditMetInfTable()

        self.AddMultiScatter()
        self.EditHeader()

        # Move the supplemental info slide to the last slide position
        slides = self.rpt.slides
        slide = slides[1]
        new_idx = len(slides)
        self.MoveSlide(slides, slide, new_idx)
        self.AddSlideNumbers()

        self.SaveReport()

    def SaveReport(self):
        """Save the report to the `/reports` directory as a pptx file.

        Returns:
            None.

        """
        print('..Saving report')
        self.rpt_name = 'Base_Testing_Report_' + self._param_name\
                        + '_' + self.name + '_' + self.today + '.pptx'

        save_dir = '\\'.join((self.path, 'reports',
                              self.name, self._param_name))
        save_path = '\\'.join((save_dir, self.rpt_name))

        if not os.path.exists(save_dir):
            os.makedirs(save_dir)
            print('..Creating directory:')
            print('....' + save_dir)

        print('....' + save_path.replace(self.path, ''))
        self.rpt.save(save_path)
